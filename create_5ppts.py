import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)  # Deep navy
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)  # Card background
ACCENT_BLUE   = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GREEN  = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
ACCENT_RED    = RGBColor(0xFF, 0x6B, 0x6B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xB0, 0xB8, 0xC4)
GRAY_DIM   = RGBColor(0x7A, 0x82, 0x8E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=14, color=WHITE, bold=False, space_before=Pt(4), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.space_before = space_before
    p.alignment = alignment
    return p

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def title_slide(prs, title, subtitle, accent_color=ACCENT_BLUE, num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    # accent bar top
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, accent_color)
    # number badge
    if num:
        add_text_box(slide, Inches(0.8), Inches(1.8), Inches(1.2), Inches(1.2), num, font_size=60, color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    # title
    add_text_box(slide, Inches(2.2), Inches(2.0), Inches(9), Inches(1.0), title, font_size=40, color=WHITE, bold=True)
    # subtitle
    add_text_box(slide, Inches(2.2), Inches(3.2), Inches(9), Inches(0.8), subtitle, font_size=18, color=GRAY_LIGHT)
    # bottom bar
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(6), Inches(0.4),
                 'BTC Market Intelligence Indicator  |  System Architecture Deep Dive',
                 font_size=11, color=GRAY_DIM)
    add_text_box(slide, Inches(9), Inches(6.5), Inches(3.5), Inches(0.4),
                 'source@rfo  |  2026', font_size=11, color=GRAY_DIM, alignment=PP_ALIGN.RIGHT)
    return slide

def section_slide(prs, title, accent_color=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_line(slide, Inches(1), Inches(3.4), Inches(2), accent_color)
    add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(1), title, font_size=36, color=WHITE, bold=True)
    return slide

def content_slide(prs, title, accent_color=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, accent_color)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6), title, font_size=24, color=WHITE, bold=True)
    return slide

def footer(slide):
    add_text_box(slide, Inches(0.6), Inches(7.0), Inches(6), Inches(0.3),
                 'BTC Market Intelligence Indicator  |  source@rfo', font_size=9, color=GRAY_DIM)

# ═══════════════════════════════════════════════
# PPT 1: 資料擷取 (Data Acquisition)
# ═══════════════════════════════════════════════
def create_ppt1():
    prs = new_prs()
    AC = ACCENT_BLUE

    # --- Slide 1: Title ---
    title_slide(prs, '資料擷取層', 'Data Acquisition Layer — 多源即時數據收集與處理', AC, '01')

    # --- Slide 2: Overview ---
    s = content_slide(prs, '資料擷取層概覽', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '什麼是資料擷取層？', 18, AC, True)
    add_para(tf, '', 6)
    add_para(tf, '資料擷取層是整個預測系統的第一步，負責從多個外部', 13, GRAY_LIGHT)
    add_para(tf, 'API 拉取原始市場數據，進行清洗、驗證後儲存為標準', 13, GRAY_LIGHT)
    add_para(tf, '化格式，供後續特徵工程使用。', 13, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '核心職責：', 14, WHITE, True)
    add_para(tf, '• 連接 Binance Futures REST API 拉取 K 線', 13, GRAY_LIGHT)
    add_para(tf, '• 連接 Coinglass API 拉取 7 種衍生品指標', 13, GRAY_LIGHT)
    add_para(tf, '• 資料驗證（正值、格式、時間連續性）', 13, GRAY_LIGHT)
    add_para(tf, '• 失敗重試 + 快取回退機制', 13, GRAY_LIGHT)
    add_para(tf, '• 新鮮度監控與異常告警', 13, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '關鍵指標', 18, AC, True)
    add_para(tf2, '', 6)
    add_para(tf2, '資料源數量          2 (Binance + Coinglass)', 12, GRAY_LIGHT)
    add_para(tf2, 'Coinglass 端點      7 個獨立 API', 12, GRAY_LIGHT)
    add_para(tf2, '拉取頻率              每小時整點', 12, GRAY_LIGHT)
    add_para(tf2, 'K 線數量              500 根 / 次（~20 天）', 12, GRAY_LIGHT)
    add_para(tf2, '重試策略              指數退避，最多 3 次', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.7))
    tf3 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), '名詞解釋', 18, ACCENT_PURPLE, True)
    add_para(tf3, '', 6)
    add_para(tf3, 'REST API: 透過 HTTP 請求取得資料的標準方式', 12, GRAY_LIGHT)
    add_para(tf3, 'OHLCV: Open/High/Low/Close/Volume 的縮寫', 12, GRAY_LIGHT)
    add_para(tf3, 'Perpetual: 無到期日的永續合約', 12, GRAY_LIGHT)
    add_para(tf3, 'Parquet: 高效率欄式儲存格式', 12, GRAY_LIGHT)
    add_para(tf3, '指數退避: 每次重試等待時間加倍（2s→4s→8s）', 12, GRAY_LIGHT)
    footer(s)

    # --- Slide 3: Binance ---
    s = content_slide(prs, 'Binance Futures API — K 線數據', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5), 'API 端點與參數', 16, AC, True)
    add_para(tf, '', 6)
    add_para(tf, '端點: GET /fapi/v1/klines', 13, WHITE)
    add_para(tf, '', 4)
    add_para(tf, '參數設定：', 13, WHITE, True)
    add_para(tf, '  symbol     = BTCUSDT          # 交易對', 12, GRAY_LIGHT)
    add_para(tf, '  interval   = 1h               # K 線週期', 12, GRAY_LIGHT)
    add_para(tf, '  limit      = 500              # 一次拉取數量', 12, GRAY_LIGHT)

    tf2 = add_text_box(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.5), '回傳欄位', 16, AC, True)
    add_para(tf2, '', 6)
    add_para(tf2, 'Open Time (ms)      — K 線開始時間戳', 12, GRAY_LIGHT)
    add_para(tf2, 'Open / High / Low / Close — 四價位', 12, GRAY_LIGHT)
    add_para(tf2, 'Volume              — 基礎資產成交量', 12, GRAY_LIGHT)
    add_para(tf2, 'Quote Volume        — 計價資產成交量 (USDT)', 12, GRAY_LIGHT)
    add_para(tf2, 'Taker Buy Volume    — 主動買方成交量', 12, GRAY_LIGHT)
    add_para(tf2, 'Trade Count         — 成交筆數', 12, GRAY_LIGHT)

    add_card(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(2.8))
    tf3 = add_text_box(s, Inches(0.8), Inches(4.2), Inches(5.3), Inches(0.5), '資料清洗步驟', 16, ACCENT_GREEN, True)
    add_para(tf3, '', 6)
    add_para(tf3, '1. 丟棄最後一根未完成 bar', 13, GRAY_LIGHT)
    add_para(tf3, '   → 避免使用不完整的 K 線導致偏差', 11, GRAY_DIM)
    add_para(tf3, '2. 價格驗證：檢查 > 0、格式合法', 13, GRAY_LIGHT)
    add_para(tf3, '   → 防止 API 回傳異常值', 11, GRAY_DIM)
    add_para(tf3, '3. 時間連續性檢查', 13, GRAY_LIGHT)
    add_para(tf3, '   → 確保沒有缺漏的 bar', 11, GRAY_DIM)
    add_para(tf3, '4. 型別轉換：字串 → float / int', 13, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.8))
    tf4 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), '為什麼用 Binance？', 16, ACCENT_ORANGE, True)
    add_para(tf4, '', 6)
    add_para(tf4, '• 全球最大加密衍生品交易所', 13, GRAY_LIGHT)
    add_para(tf4, '• BTC 永續合約交易量佔比 > 40%', 13, GRAY_LIGHT)
    add_para(tf4, '• API 穩定、文件完整、免費使用', 13, GRAY_LIGHT)
    add_para(tf4, '• 支援高頻率請求（1200 req/min）', 13, GRAY_LIGHT)
    add_para(tf4, '• Taker Buy Volume 是關鍵訂單流數據', 13, GRAY_LIGHT)
    add_para(tf4, '  → 區分主動買 vs 主動賣', 11, GRAY_DIM)
    footer(s)

    # --- Slide 4: Coinglass ---
    s = content_slide(prs, 'Coinglass API — 7 個衍生品數據端點', AC)
    # Row 1
    for i, (name, desc, detail) in enumerate([
        ('Open Interest (OI)', '未平倉合約總量', '反映市場槓桿水平\n增加 = 新資金進場\n減少 = 部位平倉'),
        ('Liquidation', '爆倉數據', '強制平倉事件統計\n大量爆倉 = 市場極端\n方向性爆倉指標'),
        ('Long/Short Ratio', '大戶多空比', '頭部帳戶多空比例\n> 1 = 大戶偏多\n< 1 = 大戶偏空'),
        ('Funding Rate', '資金費率', '多空持倉成本差異\n正 = 多方付費給空方\n負 = 空方付費給多方'),
    ]):
        x = Inches(0.5 + i * 3.1)
        add_card(s, x, Inches(1.2), Inches(2.8), Inches(2.5))
        add_text_box(s, x + Inches(0.2), Inches(1.4), Inches(2.4), Inches(0.4), name, 13, AC, True)
        add_text_box(s, x + Inches(0.2), Inches(1.9), Inches(2.4), Inches(0.3), desc, 12, WHITE)
        add_text_box(s, x + Inches(0.2), Inches(2.3), Inches(2.4), Inches(1.2), detail, 11, GRAY_LIGHT)

    # Row 2
    for i, (name, desc, detail) in enumerate([
        ('Global LS Ratio', '全網多空比', '包含所有交易所\n比單所數據更全面\n差異 = 跨所分歧'),
        ('Taker Buy/Sell Vol', '主動買賣量', '市價單方向統計\n主動買 > 賣 = 偏多\n反映即時情緒'),
        ('OI Aggregated', '全網 OI 聚合', '跨交易所 OI 匯總\n更完整的槓桿畫面\n避免單所偏差'),
    ]):
        x = Inches(0.5 + i * 3.1)
        add_card(s, x, Inches(4.0), Inches(2.8), Inches(2.5))
        add_text_box(s, x + Inches(0.2), Inches(4.2), Inches(2.4), Inches(0.4), name, 13, AC, True)
        add_text_box(s, x + Inches(0.2), Inches(4.7), Inches(2.4), Inches(0.3), desc, 12, WHITE)
        add_text_box(s, x + Inches(0.2), Inches(5.1), Inches(2.4), Inches(1.2), detail, 11, GRAY_LIGHT)

    # Why box
    add_card(s, Inches(9.8), Inches(4.0), Inches(2.8), Inches(2.5))
    tf = add_text_box(s, Inches(10.0), Inches(4.2), Inches(2.4), Inches(0.4), '速率控制', 13, ACCENT_ORANGE, True)
    add_text_box(s, Inches(10.0), Inches(4.7), Inches(2.4), Inches(1.8),
                 '每端點間隔 1 秒\n避免 rate limit\n7 端點 ≈ 7 秒完成\n\n失敗自動重試\n快取回退保障', 11, GRAY_LIGHT)
    footer(s)

    # --- Slide 5: Retry & Cache ---
    s = content_slide(prs, '容錯機制 — 重試、快取與監控', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(3.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(3.3), Inches(0.5), '指數退避重試', 16, AC, True)
    add_para(tf, '', 6)
    add_para(tf, '什麼是指數退避？', 13, WHITE, True)
    add_para(tf, '每次重試的等待時間以指數增長，', 12, GRAY_LIGHT)
    add_para(tf, '避免在 API 故障時造成更大負擔。', 12, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '第 1 次重試:  等待 2 秒', 12, ACCENT_GREEN)
    add_para(tf, '第 2 次重試:  等待 4 秒', 12, ACCENT_ORANGE)
    add_para(tf, '第 3 次重試:  等待 8 秒', 12, ACCENT_RED)
    add_para(tf, '3 次失敗後:   啟用快取回退', 12, ACCENT_RED, True)
    add_para(tf, '', 6)
    add_para(tf, '為什麼不用固定間隔？', 13, WHITE, True)
    add_para(tf, '固定間隔在 API 過載時會持續施壓，', 12, GRAY_LIGHT)
    add_para(tf, '指數退避讓伺服器有時間恢復。', 12, GRAY_LIGHT)

    add_card(s, Inches(4.7), Inches(1.2), Inches(3.8), Inches(5.5))
    tf2 = add_text_box(s, Inches(5.0), Inches(1.4), Inches(3.3), Inches(0.5), 'Parquet 快取回退', 16, ACCENT_GREEN, True)
    add_para(tf2, '', 6)
    add_para(tf2, '機制說明：', 13, WHITE, True)
    add_para(tf2, '每次 API 成功拉取後，資料同時', 12, GRAY_LIGHT)
    add_para(tf2, '儲存為 Parquet 檔案作為快取。', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, '回退流程：', 13, WHITE, True)
    add_para(tf2, '1. 嘗試 API 拉取', 12, GRAY_LIGHT)
    add_para(tf2, '2. 若 3 次重試皆失敗', 12, GRAY_LIGHT)
    add_para(tf2, '3. 自動載入上次成功的 Parquet', 12, GRAY_LIGHT)
    add_para(tf2, '4. 使用快取數據繼續運算', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, '為什麼用 Parquet？', 13, WHITE, True)
    add_para(tf2, '• 欄式儲存，讀取速度快', 12, GRAY_LIGHT)
    add_para(tf2, '• 自帶壓縮，檔案小', 12, GRAY_LIGHT)
    add_para(tf2, '• 完整保留型別資訊', 12, GRAY_LIGHT)

    add_card(s, Inches(8.9), Inches(1.2), Inches(3.8), Inches(5.5))
    tf3 = add_text_box(s, Inches(9.2), Inches(1.4), Inches(3.3), Inches(0.5), '新鮮度監控', 16, ACCENT_ORANGE, True)
    add_para(tf3, '', 6)
    add_para(tf3, '監控規則：', 13, WHITE, True)
    add_para(tf3, '每次更新時記錄最後成功時間。', 12, GRAY_LIGHT)
    add_para(tf3, '若資料超過 3 小時未更新：', 12, GRAY_LIGHT)
    add_para(tf3, '', 6)
    add_para(tf3, '→ 發送 Telegram 告警', 13, ACCENT_RED, True)
    add_para(tf3, '→ 標記系統狀態為 stale', 13, ACCENT_RED, True)
    add_para(tf3, '', 6)
    add_para(tf3, '為什麼是 3 小時？', 13, WHITE, True)
    add_para(tf3, '系統每小時更新一次，3 小時代表', 12, GRAY_LIGHT)
    add_para(tf3, '連續 3 次失敗，已超出正常波動。', 12, GRAY_LIGHT)
    add_para(tf3, '', 6)
    add_para(tf3, '告警內容包含：', 13, WHITE, True)
    add_para(tf3, '• 最後成功時間', 12, GRAY_LIGHT)
    add_para(tf3, '• 失敗的 API 端點', 12, GRAY_LIGHT)
    add_para(tf3, '• 是否已啟用快取回退', 12, GRAY_LIGHT)
    footer(s)

    # --- Slide 6: Data Flow ---
    s = content_slide(prs, '資料流程圖', AC)
    steps = [
        ('排程觸發', 'Railway cron\n每整點啟動', ACCENT_BLUE),
        ('API 請求', 'Binance + Coinglass\n共 8 個端點', ACCENT_BLUE),
        ('資料驗證', '正值檢查\n格式驗證\n時間連續性', ACCENT_GREEN),
        ('快取儲存', '成功 → Parquet\n失敗 → 讀取快取', ACCENT_ORANGE),
        ('格式轉換', 'DataFrame 標準化\n欄位命名統一', ACCENT_PURPLE),
        ('輸出', '傳遞至\n特徵工程層', ACCENT_GREEN),
    ]
    for i, (name, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 2.1)
        add_card(s, x, Inches(2.5), Inches(1.8), Inches(2.8))
        add_text_box(s, x + Inches(0.15), Inches(2.7), Inches(1.5), Inches(0.4), name, 14, color, True, PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.15), Inches(3.2), Inches(1.5), Inches(1.8), desc, 11, GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text_box(s, x + Inches(1.8), Inches(3.5), Inches(0.3), Inches(0.4), '→', 20, GRAY_DIM, alignment=PP_ALIGN.CENTER)

    add_card(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0))
    tf = add_text_box(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.5), '設計原則：每個步驟都可獨立失敗而不影響整體系統。API 層失敗有快取回退，驗證失敗有日誌記錄，全域異常捕獲確保不 crash。', 12, GRAY_LIGHT)
    footer(s)

    prs.save('PPT1_資料擷取層.pptx')
    print('PPT1 saved.')

# ═══════════════════════════════════════════════
# PPT 2: 特徵工程 (Feature Engineering)
# ═══════════════════════════════════════════════
def create_ppt2():
    prs = new_prs()
    AC = ACCENT_GREEN

    title_slide(prs, '特徵工程', 'Feature Engineering — 69 個量化特徵的設計與計算', AC, '02')

    # --- Overview ---
    s = content_slide(prs, '特徵工程概覽', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '什麼是特徵工程？', 18, AC, True)
    add_para(tf, '', 6)
    add_para(tf, '特徵工程是將原始市場數據轉換為機器學習模型可以', 13, GRAY_LIGHT)
    add_para(tf, '理解的數值化「特徵」。模型無法直接看懂 K 線，', 13, GRAY_LIGHT)
    add_para(tf, '但可以理解「過去 4 小時的收益率是 -2.3%」。', 13, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '本系統從 2 個數據源提取 69 個特徵：', 14, WHITE, True)
    add_para(tf, '', 6)
    add_para(tf, '類別                 數量    佔比', 13, AC)
    add_para(tf, '價格動量             9 個    13%', 12, GRAY_LIGHT)
    add_para(tf, '波動率               5 個     7%', 12, GRAY_LIGHT)
    add_para(tf, '成交量動態           6 個     9%', 12, GRAY_LIGHT)
    add_para(tf, 'Coinglass 衍生      49 個    71%', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '設計原則', 18, ACCENT_PURPLE, True)
    add_para(tf2, '', 6)
    add_para(tf2, '1. 全部 Backward-Looking', 14, WHITE, True)
    add_para(tf2, '   只使用當前時刻及之前的數據計算', 12, GRAY_LIGHT)
    add_para(tf2, '   絕對沒有未來函數（Look-Ahead Bias）', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, '2. 支援增量計算', 14, WHITE, True)
    add_para(tf2, '   Rolling window 與 shift 皆可即時更新', 12, GRAY_LIGHT)
    add_para(tf2, '   不需每次重算全部歷史', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, '3. 排除傳統技術指標', 14, WHITE, True)
    add_para(tf2, '   不使用 MACD / EMA / RSI / Bollinger', 12, GRAY_LIGHT)
    add_para(tf2, '   專注訂單流與衍生品數據', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, '4. 訓練與推論一致', 14, WHITE, True)
    add_para(tf2, '   同一套 feature_builder 代碼', 12, GRAY_LIGHT)
    add_para(tf2, '   消除 train-serve skew', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, '5. Z-Score 標準化', 14, WHITE, True)
    add_para(tf2, '   跨時期可比較，避免量綱差異', 12, GRAY_LIGHT)
    footer(s)

    # --- Price Momentum ---
    s = content_slide(prs, '特徵類別 1：價格動量 (Price Momentum)', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '特徵列表', 16, AC, True)
    features = [
        ('return_1h', '過去 1 小時收益率', 'close / close[-1] - 1'),
        ('return_2h', '過去 2 小時收益率', 'close / close[-2] - 1'),
        ('return_4h', '過去 4 小時收益率', 'close / close[-4] - 1'),
        ('return_8h', '過去 8 小時收益率', 'close / close[-8] - 1'),
        ('return_24h', '過去 24 小時收益率', 'close / close[-24] - 1'),
        ('return_zscore_20', '20 期收益率 Z-Score', '(ret - mean_20) / std_20'),
        ('return_zscore_50', '50 期收益率 Z-Score', '(ret - mean_50) / std_50'),
        ('close_vs_ema_50', '價格 vs 50期均線偏離', '(close - ema50) / ema50'),
        ('high_low_range', 'K 線振幅比', '(high - low) / close'),
    ]
    for fname, desc, formula in features:
        add_para(tf, '', 4)
        add_para(tf, f'{fname}', 12, WHITE, True)
        add_para(tf, f'  {desc}  |  {formula}', 11, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '名詞解釋', 16, ACCENT_PURPLE, True)
    add_para(tf2, '', 6)
    add_para(tf2, '收益率 (Return): 價格變動百分比', 12, GRAY_LIGHT)
    add_para(tf2, '  例如 BTC 從 60000 漲到 61200 = +2%', 11, GRAY_DIM)
    add_para(tf2, '', 4)
    add_para(tf2, 'Z-Score: 衡量當前值偏離平均的程度', 12, GRAY_LIGHT)
    add_para(tf2, '  = (當前值 - 均值) / 標準差', 11, GRAY_DIM)
    add_para(tf2, '  > 2 表示異常偏高，< -2 表示異常偏低', 11, GRAY_DIM)
    add_para(tf2, '', 4)
    add_para(tf2, 'EMA: 指數移動平均線，近期權重較大', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.7))
    tf3 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), '為什麼多時間框架？', 16, ACCENT_ORANGE, True)
    add_para(tf3, '', 6)
    add_para(tf3, '不同時間範圍的收益率捕捉不同信息：', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '1h: 短期微結構 — 即時波動方向', 12, GRAY_LIGHT)
    add_para(tf3, '4h: 中期趨勢 — 日內主要方向', 12, GRAY_LIGHT)
    add_para(tf3, '24h: 長期背景 — 整體牛熊環境', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '模型綜合這些尺度做出更穩健的預測。', 12, GRAY_LIGHT)
    footer(s)

    # --- Volatility ---
    s = content_slide(prs, '特徵類別 2：波動率 (Volatility)', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '波動率特徵', 16, AC, True)
    vols = [
        ('realized_vol_10b', '10 期實現波動率', 'std(returns, 10) * sqrt(24)'),
        ('realized_vol_20b', '20 期實現波動率', 'std(returns, 20) * sqrt(24)'),
        ('realized_vol_50b', '50 期實現波動率', 'std(returns, 50) * sqrt(24)'),
        ('vol_ratio', '短/長期波動率比值', 'vol_10 / vol_50'),
        ('parkinson_vol', 'Parkinson 波動率', 'sqrt(1/4nln2 * sum(ln(H/L))^2)'),
    ]
    for fname, desc, formula in vols:
        add_para(tf, '', 6)
        add_para(tf, f'{fname}', 13, WHITE, True)
        add_para(tf, f'  {desc}', 12, GRAY_LIGHT)
        add_para(tf, f'  公式: {formula}', 11, GRAY_DIM)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.8))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '名詞解釋', 16, ACCENT_PURPLE, True)
    add_para(tf2, '', 6)
    add_para(tf2, '波動率 (Volatility):', 13, WHITE, True)
    add_para(tf2, '價格變動的劇烈程度。高波動 = 大漲大跌，', 12, GRAY_LIGHT)
    add_para(tf2, '低波動 = 價格穩定。年化後可跨市場比較。', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, 'Parkinson 波動率:', 13, WHITE, True)
    add_para(tf2, '使用 High/Low 而非 Close 計算，能捕捉日內', 12, GRAY_LIGHT)
    add_para(tf2, '波動，比傳統 close-to-close 方法更精確。', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, 'Vol Ratio (短/長比):', 13, WHITE, True)
    add_para(tf2, '> 1 = 近期波動加劇（可能趨勢形成）', 12, GRAY_LIGHT)
    add_para(tf2, '< 1 = 近期波動縮小（可能盤整壓縮）', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.4))
    tf3 = add_text_box(s, Inches(7.1), Inches(4.5), Inches(5.3), Inches(0.5), '為什麼波動率重要？', 16, ACCENT_ORANGE, True)
    add_para(tf3, '', 6)
    add_para(tf3, '• 同樣 +2% 的預測，在低波動環境是「大動作」，', 12, GRAY_LIGHT)
    add_para(tf3, '  在高波動環境只是「正常波動」', 12, GRAY_LIGHT)
    add_para(tf3, '• 模型的目標是波動率調整後的收益', 12, GRAY_LIGHT)
    add_para(tf3, '• 波動率特徵幫助 Regime Detection 分類', 12, GRAY_LIGHT)
    footer(s)

    # --- Volume ---
    s = content_slide(prs, '特徵類別 3：成交量動態 (Volume Dynamics)', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.3))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5), '成交量特徵', 16, AC, True)
    for fname, desc in [
        ('volume_zscore', '成交量 Z-Score — 偏離近期平均的程度'),
        ('taker_buy_ratio', '主動買佔比 — taker_buy_vol / total_vol'),
        ('taker_buy_zscore', '主動買佔比的 Z-Score'),
    ]:
        add_para(tf, f'  {fname}: {desc}', 12, GRAY_LIGHT)

    tf2 = add_text_box(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.5), '', 16, AC, True)
    for fname, desc in [
        ('volume_trend', 'OBV 斜率 — On-Balance Volume 趨勢'),
        ('rel_volume_10', '相對成交量 10 期 — vol / sma(vol,10)'),
        ('rel_volume_50', '相對成交量 50 期 — vol / sma(vol,50)'),
    ]:
        add_para(tf2, f'  {fname}: {desc}', 12, GRAY_LIGHT)

    add_card(s, Inches(0.5), Inches(3.8), Inches(5.8), Inches(3.0))
    tf3 = add_text_box(s, Inches(0.8), Inches(4.0), Inches(5.3), Inches(0.5), '名詞解釋', 16, ACCENT_PURPLE, True)
    add_para(tf3, '', 6)
    add_para(tf3, 'Taker Buy Ratio (主動買佔比):', 13, WHITE, True)
    add_para(tf3, '市場上主動以市價「吃」賣單的成交量佔總量比例。', 12, GRAY_LIGHT)
    add_para(tf3, '> 50% = 買方更積極，< 50% = 賣方更積極', 12, GRAY_LIGHT)
    add_para(tf3, '', 6)
    add_para(tf3, 'OBV (On-Balance Volume):', 13, WHITE, True)
    add_para(tf3, '累積量價指標。上漲時加成交量，下跌時減。', 12, GRAY_LIGHT)
    add_para(tf3, 'OBV 上升 = 量能支撐上漲，下降 = 量能萎縮。', 12, GRAY_LIGHT)
    add_para(tf3, '', 6)
    add_para(tf3, '相對成交量:', 13, WHITE, True)
    add_para(tf3, '當前成交量 / 平均成交量。> 2 = 放量異常。', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(3.0))
    tf4 = add_text_box(s, Inches(7.1), Inches(4.0), Inches(5.3), Inches(0.5), '成交量在預測中的角色', 16, ACCENT_ORANGE, True)
    add_para(tf4, '', 6)
    add_para(tf4, '「成交量先於價格」— 量化交易經典法則', 13, WHITE, True)
    add_para(tf4, '', 6)
    add_para(tf4, '• 趨勢啟動前通常伴隨成交量激增', 12, GRAY_LIGHT)
    add_para(tf4, '• 放量突破 vs 縮量突破的可靠度截然不同', 12, GRAY_LIGHT)
    add_para(tf4, '• Taker buy ratio 直接反映市場即時供需', 12, GRAY_LIGHT)
    add_para(tf4, '', 6)
    add_para(tf4, '永續合約的特殊性：', 13, WHITE, True)
    add_para(tf4, '• 合約市場可做空，buy ratio 更有意義', 12, GRAY_LIGHT)
    add_para(tf4, '• 大單 taker 通常代表機構或鯨魚行為', 12, GRAY_LIGHT)
    footer(s)

    # --- Coinglass ---
    s = content_slide(prs, '特徵類別 4：Coinglass 衍生特徵（49 個）', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(3.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(3.3), Inches(0.5), 'Open Interest 系列', 14, AC, True)
    add_para(tf, '', 4)
    add_para(tf, 'oi_change_1h     — 1h OI 變化', 11, GRAY_LIGHT)
    add_para(tf, 'oi_change_4h     — 4h OI 變化', 11, GRAY_LIGHT)
    add_para(tf, 'oi_zscore_20     — 20期 OI Z-Score', 11, GRAY_LIGHT)
    add_para(tf, 'oi_zscore_50     — 50期 OI Z-Score', 11, GRAY_LIGHT)
    add_para(tf, '', 4)
    add_para(tf, 'Funding Rate 系列', 14, AC, True)
    add_para(tf, '', 4)
    add_para(tf, 'funding_rate     — 當前費率', 11, GRAY_LIGHT)
    add_para(tf, 'funding_zscore   — 費率 Z-Score', 11, GRAY_LIGHT)
    add_para(tf, 'funding_deviation — 偏離均值', 11, GRAY_LIGHT)
    add_para(tf, 'funding_cum_8h   — 8h 累積費率', 11, GRAY_LIGHT)
    add_para(tf, '', 4)
    add_para(tf, 'Long/Short Ratio 系列', 14, AC, True)
    add_para(tf, '', 4)
    add_para(tf, 'ls_ratio          — 多空比', 11, GRAY_LIGHT)
    add_para(tf, 'ls_zscore         — 多空比 Z-Score', 11, GRAY_LIGHT)
    add_para(tf, 'global_ls_ratio   — 全網多空比', 11, GRAY_LIGHT)
    add_para(tf, 'ls_divergence     — 單所 vs 全網差異', 11, GRAY_LIGHT)

    add_card(s, Inches(4.7), Inches(1.2), Inches(3.8), Inches(5.5))
    tf2 = add_text_box(s, Inches(5.0), Inches(1.4), Inches(3.3), Inches(0.5), 'Liquidation 系列', 14, AC, True)
    add_para(tf2, '', 4)
    add_para(tf2, 'liq_long_usd     — 多方爆倉量', 11, GRAY_LIGHT)
    add_para(tf2, 'liq_short_usd    — 空方爆倉量', 11, GRAY_LIGHT)
    add_para(tf2, 'liq_ratio        — 多/空爆倉比', 11, GRAY_LIGHT)
    add_para(tf2, 'liq_zscore       — 爆倉量 Z-Score', 11, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, 'Taker Volume 系列', 14, AC, True)
    add_para(tf2, '', 4)
    add_para(tf2, 'taker_buy_vol    — 主動買量', 11, GRAY_LIGHT)
    add_para(tf2, 'taker_sell_vol   — 主動賣量', 11, GRAY_LIGHT)
    add_para(tf2, 'taker_delta      — 買賣差', 11, GRAY_LIGHT)
    add_para(tf2, 'taker_zscore     — 買賣差 Z-Score', 11, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '交叉衍生', 14, AC, True)
    add_para(tf2, '', 4)
    add_para(tf2, 'oi_x_funding     — OI * Funding', 11, GRAY_LIGHT)
    add_para(tf2, 'liq_x_vol        — 爆倉 * 波動', 11, GRAY_LIGHT)
    add_para(tf2, 'funding_x_ls     — 費率 * 多空比', 11, GRAY_LIGHT)
    add_para(tf2, '... 多種交叉組合', 11, GRAY_DIM)

    add_card(s, Inches(8.9), Inches(1.2), Inches(3.8), Inches(5.5))
    tf3 = add_text_box(s, Inches(9.2), Inches(1.4), Inches(3.3), Inches(0.5), '為什麼這些數據重要？', 14, ACCENT_ORANGE, True)
    add_para(tf3, '', 6)
    add_para(tf3, 'Open Interest:', 13, WHITE, True)
    add_para(tf3, 'OI 增加 + 價格上漲 = 新多單進場', 11, GRAY_LIGHT)
    add_para(tf3, 'OI 增加 + 價格下跌 = 新空單進場', 11, GRAY_LIGHT)
    add_para(tf3, 'OI 減少 = 平倉離場，動能減弱', 11, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, 'Funding Rate:', 13, WHITE, True)
    add_para(tf3, '正費率 = 多方擁擠，有回調風險', 11, GRAY_LIGHT)
    add_para(tf3, '負費率 = 空方擁擠，有軋空風險', 11, GRAY_LIGHT)
    add_para(tf3, '極端費率 = 反轉信號', 11, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, 'Liquidation:', 13, WHITE, True)
    add_para(tf3, '大量爆倉 = 瀑布效應（清倉引發更多清倉）', 11, GRAY_LIGHT)
    add_para(tf3, '爆倉後通常出現短期反彈', 11, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, 'Long/Short Ratio:', 13, WHITE, True)
    add_para(tf3, '大戶多空比是「聰明錢」風向標', 11, GRAY_LIGHT)
    add_para(tf3, '與散戶方向相反時特別有參考價值', 11, GRAY_LIGHT)
    footer(s)

    # --- Z-Score ---
    s = content_slide(prs, 'Z-Score 標準化 — 核心處理方法', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), 'Z-Score 公式與意義', 18, AC, True)
    add_para(tf, '', 8)
    add_para(tf, 'Z = (X - μ) / σ', 24, WHITE, True)
    add_para(tf, '', 8)
    add_para(tf, 'X = 當前觀測值', 14, GRAY_LIGHT)
    add_para(tf, 'μ = 過去 N 期的平均值 (rolling mean)', 14, GRAY_LIGHT)
    add_para(tf, 'σ = 過去 N 期的標準差 (rolling std)', 14, GRAY_LIGHT)
    add_para(tf, '', 8)
    add_para(tf, '解讀：', 16, WHITE, True)
    add_para(tf, 'Z =  0    → 完全正常，等於歷史平均', 13, GRAY_LIGHT)
    add_para(tf, 'Z = +1    → 高於平均 1 個標準差', 13, ACCENT_GREEN)
    add_para(tf, 'Z = +2    → 異常偏高（95% 信賴區間外）', 13, ACCENT_ORANGE)
    add_para(tf, 'Z = +3    → 極端偏高（99.7% 區間外）', 13, ACCENT_RED)
    add_para(tf, 'Z = -2    → 異常偏低', 13, ACCENT_RED)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '為什麼要用 Z-Score？', 16, ACCENT_ORANGE, True)
    add_para(tf2, '', 6)
    add_para(tf2, '1. 消除量綱差異', 13, WHITE, True)
    add_para(tf2, '   Funding Rate 是 0.01%，OI 是十億美元級', 12, GRAY_LIGHT)
    add_para(tf2, '   Z-Score 讓所有特徵在同一尺度比較', 12, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '2. 跨時期可比', 13, WHITE, True)
    add_para(tf2, '   BTC 30000 和 100000 時的 OI 絕對值不同', 12, GRAY_LIGHT)
    add_para(tf2, '   但 Z-Score 都表示「偏離近期正常水平的程度」', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.7))
    tf3 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), 'Look-Ahead Bias 防護', 16, ACCENT_RED, True)
    add_para(tf3, '', 6)
    add_para(tf3, '什麼是 Look-Ahead Bias？', 13, WHITE, True)
    add_para(tf3, '在計算特徵時使用了「未來」的數據。', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '舉例：用全域 mean/std 做 Z-Score', 12, GRAY_LIGHT)
    add_para(tf3, '→ 等於偷看了未來的數據分布', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '正確做法：Rolling Z-Score', 13, ACCENT_GREEN, True)
    add_para(tf3, '只使用當前及之前 N 期的 mean/std', 12, GRAY_LIGHT)
    add_para(tf3, '每一個時間點只看到「過去」的資訊', 12, GRAY_LIGHT)
    footer(s)

    prs.save('PPT2_特徵工程.pptx')
    print('PPT2 saved.')

# ═══════════════════════════════════════════════
# PPT 3: 模型推論 (Model Inference)
# ═══════════════════════════════════════════════
def create_ppt3():
    prs = new_prs()
    AC = ACCENT_PURPLE

    title_slide(prs, '模型推論', 'Model Inference — Regime 偵測 + 雙目標 XGBoost 架構', AC, '03')

    # --- Overview ---
    s = content_slide(prs, '模型架構概覽', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '為什麼不用單一模型？', 18, AC, True)
    add_para(tf, '', 6)
    add_para(tf, '市場有不同的「狀態」(Regime)：', 14, WHITE, True)
    add_para(tf, '', 4)
    add_para(tf, '• 趨勢上漲 — 買方主導，動量延續', 13, ACCENT_GREEN)
    add_para(tf, '• 趨勢下跌 — 賣方主導，恐慌擴散', 13, ACCENT_RED)
    add_para(tf, '• 震盪盤整 — 沒有明確方向，假突破多', 13, ACCENT_ORANGE)
    add_para(tf, '', 6)
    add_para(tf, '不同 Regime 下，同樣的特徵組合有不同的預測力。', 13, GRAY_LIGHT)
    add_para(tf, '例如：OI 激增在趨勢中是延續信號，在震盪中是假信號。', 13, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '解決方案：Regime-Conditional 架構', 14, WHITE, True)
    add_para(tf, '', 4)
    add_para(tf, '1. 先偵測當前 Regime', 13, GRAY_LIGHT)
    add_para(tf, '2. 根據 Regime 選擇對應的模型權重', 13, GRAY_LIGHT)
    add_para(tf, '3. 混合 Global + Regime-specific 預測', 13, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '模型配置總覽', 18, AC, True)
    add_para(tf2, '', 6)
    add_para(tf2, '總模型數量: 8 個 XGBoost 模型', 16, WHITE, True)
    add_para(tf2, '', 6)
    add_para(tf2, '= 2 目標 × 4 變體', 14, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '2 個目標 (Target):', 14, WHITE, True)
    add_para(tf2, '  • up_move_vol_adj   — 預測上漲幅度', 12, ACCENT_GREEN)
    add_para(tf2, '  • down_move_vol_adj — 預測下跌幅度', 12, ACCENT_RED)
    add_para(tf2, '', 4)
    add_para(tf2, '4 個變體 (Variant):', 14, WHITE, True)
    add_para(tf2, '  • Global         — 所有數據訓練', 12, GRAY_LIGHT)
    add_para(tf2, '  • Trending Bull  — 只用牛市數據', 12, GRAY_LIGHT)
    add_para(tf2, '  • Trending Bear  — 只用熊市數據', 12, GRAY_LIGHT)
    add_para(tf2, '  • Choppy         — 只用震盪數據', 12, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, '混合權重:', 14, WHITE, True)
    add_para(tf2, '  65% Global + 35% Regime-Specific', 14, ACCENT_BLUE, True)
    add_para(tf2, '', 4)
    add_para(tf2, '→ Global 保證穩定性', 12, GRAY_LIGHT)
    add_para(tf2, '→ Regime 提供針對性調整', 12, GRAY_LIGHT)
    footer(s)

    # --- Dual Target ---
    s = content_slide(prs, '雙目標設計 (Dual Target)', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '目標定義', 18, AC, True)
    add_para(tf, '', 8)
    add_para(tf, 'Target 1: up_move_vol_adj', 14, ACCENT_GREEN, True)
    add_para(tf, '', 4)
    add_para(tf, '= max(future_high / close - 1, 0) / volatility', 13, WHITE)
    add_para(tf, '', 4)
    add_para(tf, '意義：未來 4h 內最大上漲幅度，除以波動率', 12, GRAY_LIGHT)
    add_para(tf, '→ 衡量「考慮波動後，上方還有多少空間」', 12, GRAY_LIGHT)
    add_para(tf, '', 8)
    add_para(tf, 'Target 2: down_move_vol_adj', 14, ACCENT_RED, True)
    add_para(tf, '', 4)
    add_para(tf, '= max(1 - future_low / close, 0) / volatility', 13, WHITE)
    add_para(tf, '', 4)
    add_para(tf, '意義：未來 4h 內最大下跌幅度，除以波動率', 12, GRAY_LIGHT)
    add_para(tf, '→ 衡量「考慮波動後，下方有多少風險」', 12, GRAY_LIGHT)
    add_para(tf, '', 8)
    add_para(tf, '方向判定:', 14, WHITE, True)
    add_para(tf, 'strength = up_pred - down_pred', 13, ACCENT_BLUE)
    add_para(tf, '|strength| < deadzone(0.15) → NEUTRAL', 12, GRAY_LIGHT)
    add_para(tf, 'strength > 0 → UP', 12, ACCENT_GREEN)
    add_para(tf, 'strength < 0 → DOWN', 12, ACCENT_RED)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '為什麼用雙目標？', 16, ACCENT_ORANGE, True)
    add_para(tf2, '', 6)
    add_para(tf2, '傳統單目標: 直接預測收益率 (e.g., +1.5%)', 13, WHITE)
    add_para(tf2, '→ 問題：上漲 2% 後回跌到 +0.5%，收益率', 12, GRAY_LIGHT)
    add_para(tf2, '   只顯示 +0.5%，但上方確實有 2% 空間', 12, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '雙目標: 分別預測上方和下方的「最大移動」', 13, WHITE)
    add_para(tf2, '→ 能捕捉上下兩個方向的潛在空間', 12, GRAY_LIGHT)
    add_para(tf2, '→ 非對稱性: up_pred >> down_pred = 偏多', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.7))
    tf3 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), '為什麼除以波動率？', 16, ACCENT_ORANGE, True)
    add_para(tf3, '', 6)
    add_para(tf3, '波動率調整 (Volatility Adjustment):', 13, WHITE, True)
    add_para(tf3, '', 4)
    add_para(tf3, '牛市時 BTC 日均波動 5%，+3% 是「正常」', 12, GRAY_LIGHT)
    add_para(tf3, '盤整時 BTC 日均波動 1%，+3% 是「巨大」', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '除以波動率後：', 13, WHITE, True)
    add_para(tf3, '→ 不同時期的目標值在同一尺度', 12, GRAY_LIGHT)
    add_para(tf3, '→ 模型不會偏向高波動時期', 12, GRAY_LIGHT)
    add_para(tf3, '→ 預測更穩健、更泛化', 12, GRAY_LIGHT)
    footer(s)

    # --- Regime ---
    s = content_slide(prs, 'Regime 偵測 — 市場狀態分類', AC)
    regimes = [
        ('TRENDING_BULL', '趨勢上漲', 'vol > 60th percentile\nAND ret_24h > +0.5%', '高波動 + 明確上漲\n動量延續概率高', ACCENT_GREEN),
        ('TRENDING_BEAR', '趨勢下跌', 'vol > 60th percentile\nAND ret_24h < -0.5%', '高波動 + 明確下跌\n恐慌擴散概率高', ACCENT_RED),
        ('CHOPPY', '震盪盤整', '以上條件皆不符\n（預設狀態）', '低波動或方向不明\n假突破頻繁', ACCENT_ORANGE),
        ('WARMUP', '暖機期', '前 168 根 bar\n（= 7 天）', '數據不足以計算\n不輸出預測', GRAY_DIM),
    ]
    for i, (name, cn, rule, desc, color) in enumerate(regimes):
        x = Inches(0.5 + i * 3.1)
        add_card(s, x, Inches(1.2), Inches(2.8), Inches(5.5))
        add_text_box(s, x + Inches(0.2), Inches(1.4), Inches(2.4), Inches(0.4), name, 14, color, True, PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.2), Inches(1.9), Inches(2.4), Inches(0.3), cn, 13, WHITE, alignment=PP_ALIGN.CENTER)
        add_accent_line(s, x + Inches(0.4), Inches(2.4), Inches(2.0), color)
        add_text_box(s, x + Inches(0.2), Inches(2.6), Inches(2.4), Inches(0.4), '判定條件：', 12, GRAY_DIM)
        add_text_box(s, x + Inches(0.2), Inches(3.0), Inches(2.4), Inches(1.2), rule, 11, GRAY_LIGHT)
        add_text_box(s, x + Inches(0.2), Inches(4.3), Inches(2.4), Inches(0.4), '市場特徵：', 12, GRAY_DIM)
        add_text_box(s, x + Inches(0.2), Inches(4.7), Inches(2.4), Inches(1.2), desc, 11, GRAY_LIGHT)
    footer(s)

    # --- XGBoost ---
    s = content_slide(prs, 'XGBoost 模型規格與超參數', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '什麼是 XGBoost？', 18, AC, True)
    add_para(tf, '', 6)
    add_para(tf, 'XGBoost = eXtreme Gradient Boosting', 14, WHITE, True)
    add_para(tf, '', 4)
    add_para(tf, '一種基於決策樹的集成學習方法。', 13, GRAY_LIGHT)
    add_para(tf, '核心思想：訓練多棵「弱」決策樹，每棵新樹', 13, GRAY_LIGHT)
    add_para(tf, '專門修正前面所有樹的錯誤，最終組合成強模型。', 13, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '為什麼選擇 XGBoost？', 14, WHITE, True)
    add_para(tf, '', 4)
    add_para(tf, '• 表格數據（tabular data）表現最佳', 12, GRAY_LIGHT)
    add_para(tf, '• 訓練速度快，適合快速迭代', 12, GRAY_LIGHT)
    add_para(tf, '• 內建正則化，防止過擬合', 12, GRAY_LIGHT)
    add_para(tf, '• 支援特徵重要性分析', 12, GRAY_LIGHT)
    add_para(tf, '• 量化金融領域驗證充分', 12, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '對比深度學習：', 14, WHITE, True)
    add_para(tf, '69 個特徵 + 數千筆資料 → XGBoost 更合適', 12, GRAY_LIGHT)
    add_para(tf, '深度學習需要更多數據才能發揮優勢', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '超參數設定與解釋', 18, AC, True)
    params = [
        ('n_estimators = 600', '決策樹數量。600 棵樹依序糾錯，\n足夠學習複雜模式，配合 early stopping 防止過多。'),
        ('max_depth = 5', '每棵樹最大深度。限制為 5 層避免過擬合，\n每棵樹只捕捉中等複雜度的交互作用。'),
        ('learning_rate = 0.02', '學習率（步長）。每棵新樹只修正 2% 的殘差，\n小步長 + 多棵樹 = 更穩健的模型。'),
        ('subsample = 0.8', '每棵樹使用 80% 的訓練數據。\n隨機抽樣增加多樣性，減少過擬合。'),
        ('Early Stopping', '驗證集誤差不再下降時自動停止訓練。\n避免多餘的樹只學到噪聲。'),
        ('L1/L2 正則化', 'L1 (Lasso) 迫使不重要特徵權重歸零。\nL2 (Ridge) 懲罰過大的權重，平滑預測。'),
    ]
    for name, desc in params:
        add_para(tf2, '', 4)
        add_para(tf2, name, 13, ACCENT_BLUE, True)
        for line in desc.split('\n'):
            add_para(tf2, f'  {line}', 11, GRAY_LIGHT)
    footer(s)

    # --- Walk-Forward CV ---
    s = content_slide(prs, 'Walk-Forward Cross-Validation', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5), '什麼是 Walk-Forward CV？', 18, AC, True)
    add_para(tf, '', 6)
    add_para(tf, '時序資料不能使用傳統的隨機 K-Fold 交叉驗證（會造成未來數據洩漏）。', 14, GRAY_LIGHT)
    add_para(tf, 'Walk-Forward 嚴格保持時間順序：訓練集永遠在測試集「之前」。', 14, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, 'Fold 1:  [■■■■ Train ■■■■][▓▓ Test ▓▓]                          ', 13, WHITE)
    add_para(tf, 'Fold 2:  [■■■■■■ Train ■■■■■■][▓▓ Test ▓▓]                    ', 13, WHITE)
    add_para(tf, 'Fold 3:  [■■■■■■■■ Train ■■■■■■■■][▓▓ Test ▓▓]              ', 13, WHITE)
    add_para(tf, 'Fold 4:  [■■■■■■■■■■ Train ■■■■■■■■■■][▓▓ Test ▓▓]        ', 13, WHITE)
    add_para(tf, 'Fold 5:  [■■■■■■■■■■■■ Train ■■■■■■■■■■■■][▓▓ Test ▓▓]  ', 13, WHITE)

    add_card(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(2.7))
    tf2 = add_text_box(s, Inches(0.8), Inches(4.2), Inches(5.3), Inches(0.5), '為什麼不用隨機 K-Fold？', 16, ACCENT_RED, True)
    add_para(tf2, '', 6)
    add_para(tf2, '隨機分割會將未來的數據放入訓練集：', 13, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '• 模型「看過」未來的市場狀態', 12, GRAY_LIGHT)
    add_para(tf2, '• 驗證指標虛高（假的好成績）', 12, GRAY_LIGHT)
    add_para(tf2, '• 實際部署後表現遠不如預期', 12, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '金融時序數據有自相關性，隨機分割', 13, GRAY_LIGHT)
    add_para(tf2, '等於把答案洩漏給模型。', 13, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.7))
    tf3 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), 'Walk-Forward 的優勢', 16, ACCENT_GREEN, True)
    add_para(tf3, '', 6)
    add_para(tf3, '• 模擬真實部署場景', 13, GRAY_LIGHT)
    add_para(tf3, '  → 訓練到某個時點，預測下一段', 12, GRAY_DIM)
    add_para(tf3, '• 檢驗模型在不同時期的穩定性', 13, GRAY_LIGHT)
    add_para(tf3, '  → 5 個 fold 代表 5 個不同時期', 12, GRAY_DIM)
    add_para(tf3, '• 訓練集逐步擴大', 13, GRAY_LIGHT)
    add_para(tf3, '  → 模型能利用越來越多歷史知識', 12, GRAY_DIM)
    add_para(tf3, '• 評估指標真實可靠', 13, GRAY_LIGHT)
    add_para(tf3, '  → IC 0.23~0.26 是真正的 OOS 表現', 12, GRAY_DIM)
    footer(s)

    prs.save('PPT3_模型推論.pptx')
    print('PPT3 saved.')

# ═══════════════════════════════════════════════
# PPT 4: 信號生成 (Signal Generation)
# ═══════════════════════════════════════════════
def create_ppt4():
    prs = new_prs()
    AC = ACCENT_ORANGE

    title_slide(prs, '信號生成', 'Signal Generation — 方向判定、信心校準與強度分級', AC, '04')

    # --- Overview ---
    s = content_slide(prs, '信號生成流程概覽', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.0))
    steps = [
        ('模型原始輸出', 'up_pred\ndown_pred', ACCENT_PURPLE),
        ('方向判定', 'UP / DOWN\n/ NEUTRAL', AC),
        ('信心校準', '0 ~ 100\n百分位排名', ACCENT_BLUE),
        ('強度分級', 'Strong / Moderate\n/ Weak', ACCENT_GREEN),
        ('最終信號', '圖表標記\n+ 文字描述', ACCENT_RED),
    ]
    for i, (name, desc, color) in enumerate(steps):
        x = Inches(0.6 + i * 2.45)
        add_text_box(s, x, Inches(1.4), Inches(2.0), Inches(0.4), name, 14, color, True, PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(1.9), Inches(2.0), Inches(1.0), desc, 11, GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text_box(s, x + Inches(2.0), Inches(1.7), Inches(0.4), Inches(0.4), '→', 18, GRAY_DIM, alignment=PP_ALIGN.CENTER)

    # Direction
    add_card(s, Inches(0.5), Inches(3.5), Inches(3.8), Inches(3.3))
    tf = add_text_box(s, Inches(0.8), Inches(3.7), Inches(3.3), Inches(0.5), '步驟 1: 方向判定', 16, AC, True)
    add_para(tf, '', 6)
    add_para(tf, 'strength = up_pred - down_pred', 13, WHITE, True)
    add_para(tf, '', 4)
    add_para(tf, '這個差值代表「上方空間 vs 下方風險」', 12, GRAY_LIGHT)
    add_para(tf, '的不對稱程度。', 12, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, 'Deadzone 機制:', 13, WHITE, True)
    add_para(tf, '|strength| < 0.15 → NEUTRAL', 12, ACCENT_ORANGE)
    add_para(tf, '  避免在上下空間差不多時強行選邊', 11, GRAY_DIM)
    add_para(tf, '', 4)
    add_para(tf, 'strength > 0.15 → UP', 12, ACCENT_GREEN)
    add_para(tf, 'strength < -0.15 → DOWN', 12, ACCENT_RED)

    # Confidence
    add_card(s, Inches(4.7), Inches(3.5), Inches(3.8), Inches(3.3))
    tf2 = add_text_box(s, Inches(5.0), Inches(3.7), Inches(3.3), Inches(0.5), '步驟 2: 信心校準', 16, ACCENT_BLUE, True)
    add_para(tf2, '', 6)
    add_para(tf2, '兩階段計算:', 13, WHITE, True)
    add_para(tf2, '', 4)
    add_para(tf2, '1. Magnitude Score', 13, ACCENT_BLUE, True)
    add_para(tf2, '   |strength| 在歷史分布中的百分位', 12, GRAY_LIGHT)
    add_para(tf2, '   → 衡量「這次預測有多強烈」', 11, GRAY_DIM)
    add_para(tf2, '', 4)
    add_para(tf2, '2. Regime Score', 13, ACCENT_BLUE, True)
    add_para(tf2, '   該 regime 下模型近期 IC 表現', 12, GRAY_LIGHT)
    add_para(tf2, '   → 衡量「當前模型有多可信」', 11, GRAY_DIM)
    add_para(tf2, '', 4)
    add_para(tf2, 'Confidence = mag × regime', 14, WHITE, True)
    add_para(tf2, '範圍: 0 ~ 100', 12, GRAY_LIGHT)

    # Strength
    add_card(s, Inches(8.9), Inches(3.5), Inches(3.8), Inches(3.3))
    tf3 = add_text_box(s, Inches(9.2), Inches(3.7), Inches(3.3), Inches(0.5), '步驟 3: 強度分級', 16, ACCENT_GREEN, True)
    add_para(tf3, '', 6)
    add_para(tf3, '門檻設定:', 13, WHITE, True)
    add_para(tf3, '', 4)
    add_para(tf3, 'Strong    ≥ 90', 16, ACCENT_GREEN, True)
    add_para(tf3, '  最高信心，約佔 ~10% 信號', 11, GRAY_DIM)
    add_para(tf3, '  → 唯一會在圖表上標記的信號', 11, GRAY_DIM)
    add_para(tf3, '', 4)
    add_para(tf3, 'Moderate  ≥ 65', 14, ACCENT_ORANGE, True)
    add_para(tf3, '  中等信心，僅文字顯示', 11, GRAY_DIM)
    add_para(tf3, '', 4)
    add_para(tf3, 'Weak      < 65', 14, GRAY_DIM, True)
    add_para(tf3, '  低信心，一般不顯示', 11, GRAY_DIM)
    add_para(tf3, '', 6)
    add_para(tf3, '設計理念: 寧缺勿濫', 13, WHITE, True)
    add_para(tf3, 'Strong 門檻設高，減少噪音', 12, GRAY_LIGHT)
    footer(s)

    # --- Calibration ---
    s = content_slide(prs, '校準驗證 (Calibration)', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '什麼是 Calibration？', 18, AC, True)
    add_para(tf, '', 6)
    add_para(tf, '校準 = 預測的信心分數是否「名副其實」', 14, WHITE, True)
    add_para(tf, '', 6)
    add_para(tf, '好的校準：', 14, ACCENT_GREEN, True)
    add_para(tf, '• Strong 信號的實際收益最大', 12, GRAY_LIGHT)
    add_para(tf, '• Moderate 信號次之', 12, GRAY_LIGHT)
    add_para(tf, '• Weak 信號最小', 12, GRAY_LIGHT)
    add_para(tf, '→ 預測越強 → 實際收益越高（單調遞增）', 12, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '壞的校準：', 14, ACCENT_RED, True)
    add_para(tf, '• Strong 信號的實際收益反而比 Moderate 低', 12, GRAY_LIGHT)
    add_para(tf, '→ 信心分數失去意義，不可信', 12, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '驗證方法：', 14, WHITE, True)
    add_para(tf, '將信心分數分成 5~10 個 bin，', 12, GRAY_LIGHT)
    add_para(tf, '計算每個 bin 的平均實際收益，', 12, GRAY_LIGHT)
    add_para(tf, '檢查是否嚴格單調遞增。', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '本系統的校準結果', 16, ACCENT_GREEN, True)
    add_para(tf2, '', 6)
    add_para(tf2, '單調性驗證: 通過 ✓', 14, ACCENT_GREEN, True)
    add_para(tf2, '', 4)
    add_para(tf2, '預測越強 → 實際收益越高', 13, GRAY_LIGHT)
    add_para(tf2, '', 6)
    add_para(tf2, 'Strong 信號佔比: ~10%', 14, WHITE, True)
    add_para(tf2, '', 4)
    add_para(tf2, '高門檻確保信號品質。', 13, GRAY_LIGHT)
    add_para(tf2, '如果 Strong 佔比太高（如 30%+），', 12, GRAY_LIGHT)
    add_para(tf2, '代表門檻太低，信號氾濫。', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.7))
    tf3 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), '名詞解釋', 16, ACCENT_PURPLE, True)
    add_para(tf3, '', 6)
    add_para(tf3, 'Deadzone: 中性區間，避免在不確定時', 12, GRAY_LIGHT)
    add_para(tf3, '  強行判斷方向。0.15 是經驗值。', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '百分位 (Percentile): 在所有歷史值中的', 12, GRAY_LIGHT)
    add_para(tf3, '  排名位置。90th = 超過 90% 的歷史值。', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, 'IC (Information Coefficient): 預測值與', 12, GRAY_LIGHT)
    add_para(tf3, '  實際收益的排序相關。IC=0.25 在量化', 12, GRAY_LIGHT)
    add_para(tf3, '  領域是非常好的表現。', 12, GRAY_LIGHT)
    footer(s)

    prs.save('PPT4_信號生成.pptx')
    print('PPT4 saved.')

# ═══════════════════════════════════════════════
# PPT 5: 輸出推送 (Output & Delivery)
# ═══════════════════════════════════════════════
def create_ppt5():
    prs = new_prs()
    AC = ACCENT_RED

    title_slide(prs, '輸出推送', 'Output & Delivery — 圖表渲染、API 端點與 Telegram 整合', AC, '05')

    # --- Chart Design ---
    s = content_slide(prs, '圖表渲染設計', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(3.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(3.3), Inches(0.5), 'Panel 1: Confidence Heatmap', 14, ACCENT_PURPLE, True)
    add_para(tf, '', 4)
    add_para(tf, '頂部紫色色帶', 12, GRAY_LIGHT)
    add_para(tf, '顏色深淺 = confidence 分數 (0~100)', 12, GRAY_LIGHT)
    add_para(tf, '越深 = 模型越有信心', 12, GRAY_LIGHT)
    add_para(tf, '', 6)
    add_para(tf, '作用：讓觀察者一眼看出哪些時段', 12, GRAY_LIGHT)
    add_para(tf, '模型特別有信心或特別不確定。', 12, GRAY_LIGHT)
    add_para(tf, '', 8)
    add_para(tf, 'Panel 2: K 線 + 信號三角形', 14, ACCENT_BLUE, True)
    add_para(tf, '', 4)
    add_para(tf, '標準 OHLC K 線圖', 12, GRAY_LIGHT)
    add_para(tf, '僅顯示 Strong 信號的三角形標記', 12, GRAY_LIGHT)
    add_para(tf, '', 4)
    add_para(tf, '▲ 深綠三角 = Strong UP 信號', 12, ACCENT_GREEN)
    add_para(tf, '▼ 深紅三角 = Strong DOWN 信號', 12, ACCENT_RED)
    add_para(tf, '', 4)
    add_para(tf, '三角形透明度隨 confidence 調整', 12, GRAY_LIGHT)
    add_para(tf, '→ 信心越高，三角形越不透明', 12, GRAY_LIGHT)

    add_card(s, Inches(4.7), Inches(1.2), Inches(3.8), Inches(5.5))
    tf2 = add_text_box(s, Inches(5.0), Inches(1.4), Inches(3.3), Inches(0.5), 'Panel 3: Bull/Bear Power', 14, ACCENT_GREEN, True)
    add_para(tf2, '', 4)
    add_para(tf2, '底部柱狀圖 (-1 ~ +1)', 12, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '綜合多個衍生品指標：', 12, GRAY_LIGHT)
    add_para(tf2, '• Open Interest 變化', 11, GRAY_LIGHT)
    add_para(tf2, '• Funding Rate 方向', 11, GRAY_LIGHT)
    add_para(tf2, '• Long/Short Ratio', 11, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '綠色 = 多方優勢', 12, ACCENT_GREEN)
    add_para(tf2, '紅色 = 空方優勢', 12, ACCENT_RED)
    add_para(tf2, '', 8)
    add_para(tf2, '圖表時間範圍', 14, WHITE, True)
    add_para(tf2, '', 4)
    add_para(tf2, '顯示最近 200 根 1H bar', 12, GRAY_LIGHT)
    add_para(tf2, '≈ 8.3 天的完整歷史', 12, GRAY_LIGHT)
    add_para(tf2, '', 4)
    add_para(tf2, '為什麼 200 根？', 12, ACCENT_ORANGE)
    add_para(tf2, '太少看不到趨勢全貌', 11, GRAY_LIGHT)
    add_para(tf2, '太多圖表過於密集難以閱讀', 11, GRAY_LIGHT)
    add_para(tf2, '200 根是可讀性與信息量的平衡', 11, GRAY_LIGHT)

    add_card(s, Inches(8.9), Inches(1.2), Inches(3.8), Inches(5.5))
    tf3 = add_text_box(s, Inches(9.2), Inches(1.4), Inches(3.3), Inches(0.5), '圖表設計理念', 14, ACCENT_ORANGE, True)
    add_para(tf3, '', 6)
    add_para(tf3, '最少資訊原則：', 13, WHITE, True)
    add_para(tf3, '只標記 Strong 信號', 12, GRAY_LIGHT)
    add_para(tf3, '→ 減少視覺噪音', 12, GRAY_LIGHT)
    add_para(tf3, '→ 看到三角形 = 值得關注', 12, GRAY_LIGHT)
    add_para(tf3, '', 6)
    add_para(tf3, '技術實現：', 13, WHITE, True)
    add_para(tf3, '• Matplotlib 自訂渲染', 12, GRAY_LIGHT)
    add_para(tf3, '• In-memory PNG 生成', 12, GRAY_LIGHT)
    add_para(tf3, '  → 不寫入磁碟，直接返回', 11, GRAY_DIM)
    add_para(tf3, '• 暗色主題（深色背景）', 12, GRAY_LIGHT)
    add_para(tf3, '  → 更適合交易員使用環境', 11, GRAY_DIM)
    add_para(tf3, '• 右下角 source@rfo 浮水印', 12, GRAY_LIGHT)
    add_para(tf3, '', 6)
    add_para(tf3, '更新頻率：', 13, WHITE, True)
    add_para(tf3, '每小時自動重新渲染', 12, GRAY_LIGHT)
    add_para(tf3, '包含最新的預測結果', 12, GRAY_LIGHT)
    footer(s)

    # --- Telegram ---
    s = content_slide(prs, 'Telegram Bot 整合', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), '自動推送機制', 16, AC, True)
    add_para(tf, '', 6)
    add_para(tf, '每小時定時推送（整點觸發）：', 13, WHITE, True)
    add_para(tf, '', 4)
    add_para(tf, '1. 最新圖表 PNG（含 K 線、信號、Power 柱狀圖）', 12, GRAY_LIGHT)
    add_para(tf, '2. 預測摘要文字：', 12, GRAY_LIGHT)
    add_para(tf, '   「未來 4h 預測偏多，預估漲幅 X%，信心 Y 分」', 12, ACCENT_BLUE)
    add_para(tf, '3. 當前 Regime 狀態', 12, GRAY_LIGHT)
    add_para(tf, '', 4)
    add_para(tf, 'Strong 信號額外告警：', 13, ACCENT_RED, True)
    add_para(tf, '當出現 Strong 信號時，額外發送獨立告警通知', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), '主動查詢指令', 16, ACCENT_BLUE, True)
    add_para(tf2, '', 6)
    add_para(tf2, '/chart    — 取得最新圖表', 13, WHITE)
    add_para(tf2, '  立即回傳當前最新的預測圖表 PNG', 11, GRAY_DIM)
    add_para(tf2, '', 4)
    add_para(tf2, '/status   — 查看系統狀態', 13, WHITE)
    add_para(tf2, '  最後更新時間、模型狀態、資料新鮮度', 11, GRAY_DIM)
    add_para(tf2, '', 4)
    add_para(tf2, '/help     — 使用說明', 13, WHITE)
    add_para(tf2, '  列出所有可用指令與說明', 11, GRAY_DIM)

    add_card(s, Inches(0.5), Inches(4.0), Inches(5.8), Inches(2.7))
    tf3 = add_text_box(s, Inches(0.8), Inches(4.2), Inches(5.3), Inches(0.5), '安全機制', 16, ACCENT_ORANGE, True)
    add_para(tf3, '', 6)
    add_para(tf3, 'Webhook 模式（非 Polling）:', 13, WHITE, True)
    add_para(tf3, '• Telegram 主動推送到我們的 /webhook 端點', 12, GRAY_LIGHT)
    add_para(tf3, '• 比 polling 低延遲、省資源', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '白名單機制 (ALLOWED_CHAT_IDS):', 13, WHITE, True)
    add_para(tf3, '• 只允許授權的 Chat ID 使用指令', 12, GRAY_LIGHT)
    add_para(tf3, '• 防止未授權用戶存取預測數據', 12, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '發送失敗處理:', 13, WHITE, True)
    add_para(tf3, '• Log 記錄錯誤，不中斷主流程', 12, GRAY_LIGHT)
    add_para(tf3, '• Telegram 不可用不影響模型運行', 12, GRAY_LIGHT)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.7))
    tf4 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), '名詞解釋', 16, ACCENT_PURPLE, True)
    add_para(tf4, '', 6)
    add_para(tf4, 'Webhook: 伺服器主動向指定 URL 推送事件，', 12, GRAY_LIGHT)
    add_para(tf4, '  而非客戶端輪詢。更即時、更省資源。', 12, GRAY_LIGHT)
    add_para(tf4, '', 4)
    add_para(tf4, 'Polling: 客戶端定期詢問「有新消息嗎？」', 12, GRAY_LIGHT)
    add_para(tf4, '  延遲高（取決於輪詢間隔），浪費 API 呼叫。', 12, GRAY_LIGHT)
    add_para(tf4, '', 4)
    add_para(tf4, 'Chat ID: Telegram 中每個聊天室的唯一識別碼。', 12, GRAY_LIGHT)
    add_para(tf4, '  白名單只允許特定 Chat ID 互動。', 12, GRAY_LIGHT)
    add_para(tf4, '', 4)
    add_para(tf4, 'In-memory PNG: 圖表直接在記憶體中生成，', 12, GRAY_LIGHT)
    add_para(tf4, '  不寫檔案到磁碟，更快且更安全。', 12, GRAY_LIGHT)
    footer(s)

    # --- API ---
    s = content_slide(prs, 'API 端點與部署架構', AC)
    add_card(s, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5))
    tf = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.3), Inches(0.5), 'Flask API 端點', 16, AC, True)
    add_para(tf, '', 6)
    endpoints = [
        ('GET /', '最新圖表 PNG', '直接返回圖片，可在瀏覽器中查看\n無需登入，方便快速檢查'),
        ('GET /health', '健康檢查 JSON', '返回系統狀態、最後更新時間\n用於 Railway 健康監控'),
        ('GET /json', '最新預測數據', '返回 JSON 格式的完整預測\n包含 direction, confidence, regime 等'),
        ('POST /webhook', 'Telegram 指令接收', '接收 Telegram Bot 的 webhook\n處理 /chart, /status, /help 等指令'),
    ]
    for ep, name, desc in endpoints:
        add_para(tf, ep, 14, ACCENT_BLUE, True)
        add_para(tf, f'  {name}', 12, WHITE)
        for line in desc.split('\n'):
            add_para(tf, f'  {line}', 11, GRAY_DIM)
        add_para(tf, '', 4)

    add_card(s, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.5))
    tf2 = add_text_box(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5), 'Railway 部署架構', 16, AC, True)
    add_para(tf2, '', 6)
    add_para(tf2, 'Docker: Python 3.11-slim 基底', 13, WHITE)
    add_para(tf2, 'Server: Gunicorn (1 worker, 120s timeout)', 13, WHITE)
    add_para(tf2, 'Scheduler: APScheduler 背景排程', 13, WHITE)
    add_para(tf2, 'Deploy: GitHub push → 自動部署 (CI/CD)', 13, WHITE)
    add_para(tf2, 'Config: Railway Dashboard 環境變數', 13, WHITE)

    add_card(s, Inches(6.8), Inches(4.0), Inches(5.8), Inches(2.7))
    tf3 = add_text_box(s, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.5), 'Railway 多服務架構', 16, ACCENT_PURPLE, True)
    add_para(tf3, '', 6)
    add_para(tf3, 'Service 1: 主 Bot', 13, ACCENT_BLUE, True)
    add_para(tf3, '  TradingView + OKX WebSocket', 11, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, 'Service 2: Market Data', 13, ACCENT_GREEN, True)
    add_para(tf3, '  多交易所訂單流收集', 11, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, 'Service 3: Indicator (本系統)', 13, AC, True)
    add_para(tf3, '  獨立運行，不依賴其他服務', 11, GRAY_LIGHT)
    add_para(tf3, '', 4)
    add_para(tf3, '共享 MySQL (Railway 內網)', 13, WHITE, True)
    add_para(tf3, '  三個服務透過內網連接同一資料庫', 11, GRAY_LIGHT)
    footer(s)

    # --- Fault Tolerance ---
    s = content_slide(prs, '容錯與優雅降級設計', AC)
    faults = [
        ('API 拉取失敗', '指數退避重試 3 次\n→ 失敗後回退至 Parquet 快取\n→ 使用上次成功的數據繼續運算', ACCENT_BLUE),
        ('模型載入失敗', '主模型（Regime 條件模型）失敗\n→ 回退至 Legacy 單模型\n→ 預測品質下降但不中斷', ACCENT_PURPLE),
        ('Telegram 發送失敗', 'Log 記錄錯誤詳情\n→ 不中斷主流程\n→ 下次週期自動恢復', ACCENT_ORANGE),
        ('未預期異常', '全域 try/except 捕獲\n→ 狀態標記為 error\n→ 系統不 crash，等待下次週期', ACCENT_RED),
    ]
    for i, (name, desc, color) in enumerate(faults):
        x = Inches(0.5 + i * 3.1)
        add_card(s, x, Inches(1.2), Inches(2.8), Inches(3.5))
        add_text_box(s, x + Inches(0.2), Inches(1.4), Inches(2.4), Inches(0.4), name, 14, color, True, PP_ALIGN.CENTER)
        add_accent_line(s, x + Inches(0.4), Inches(1.9), Inches(2.0), color)
        add_text_box(s, x + Inches(0.2), Inches(2.1), Inches(2.4), Inches(2.2), desc, 11, GRAY_LIGHT)

    add_card(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5))
    tf = add_text_box(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.5), '設計哲學', 16, WHITE, True)
    add_para(tf, '', 4)
    add_para(tf, '「永不 crash」— 無論發生什麼異常，系統都應該保持運行。寧可用稍舊的快取數據輸出預測，也不要完全停止。', 13, GRAY_LIGHT)
    add_para(tf, '因為在加密市場中，即使是不完美的預測也比沒有預測好。24/7 可用性是雲端部署的核心要求。', 13, GRAY_LIGHT)
    footer(s)

    prs.save('PPT5_輸出推送.pptx')
    print('PPT5 saved.')


# ═══════════════════════════════════════════════
# Generate all 5 PPTs
# ═══════════════════════════════════════════════
if __name__ == '__main__':
    create_ppt1()
    create_ppt2()
    create_ppt3()
    create_ppt4()
    create_ppt5()
    print('\nAll 5 PPTs generated successfully!')
