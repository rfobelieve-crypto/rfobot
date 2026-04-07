"""
BTC Market Intelligence Indicator — 系統架構流程簡報 v7
深色專業主題 | 全中文 | 詳細架構與數據流程

Usage: python create_architecture_ppt.py
Output: BTC_系統架構_v7.pptx
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── 色彩系統 ───
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)
BG_CARD2   = RGBColor(0x1C, 0x22, 0x2B)
ACCENT_BLUE   = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GREEN  = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
ACCENT_RED    = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_CYAN   = RGBColor(0x56, 0xD4, 0xE0)
ACCENT_YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xB0, 0xB8, 0xC4)
GRAY_DIM   = RGBColor(0x7A, 0x82, 0x8E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_ZH = 'Microsoft JhengHei'  # 微軟正黑體
FONT_EN = 'Calibri'

# ─── 基礎工具 ───

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, size=14, color=WHITE,
       bold=False, align=PP_ALIGN.LEFT, font=FONT_ZH):
    """Add a text box and return its text_frame."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf

def ap(tf, text, size=14, color=WHITE, bold=False, sp=Pt(3), align=PP_ALIGN.LEFT):
    """Add paragraph to text_frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_ZH
    p.space_before = sp
    p.alignment = align
    return p

def card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def arrow_shape(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def footer(slide):
    tb(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
       'BTC Market Intelligence Indicator  |  System Architecture v7', 9, GRAY_DIM, font=FONT_EN)
    tb(slide, Inches(9.5), Inches(7.05), Inches(3.5), Inches(0.3),
       'source@rfo  |  2026-04', 9, GRAY_DIM, align=PP_ALIGN.RIGHT, font=FONT_EN)

# ─── 版面模板 ───

def make_title_slide(prs, title, subtitle, tags):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    line(s, Inches(0), Inches(0), SLIDE_W, ACCENT_BLUE)
    line(s, Inches(0), Inches(7.4), SLIDE_W, ACCENT_BLUE)
    tb(s, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
       title, 48, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(1), Inches(3.4), Inches(11.3), Inches(0.8),
       subtitle, 20, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    tb(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
       tags, 14, ACCENT_BLUE, align=PP_ALIGN.CENTER, font=FONT_EN)
    footer(s)
    return s

def make_section(prs, num, title, subtitle, color=ACCENT_BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    line(s, Inches(0), Inches(0), SLIDE_W, color)
    tb(s, Inches(0.8), Inches(1.8), Inches(1.8), Inches(1.8),
       num, 80, color, True, PP_ALIGN.CENTER, FONT_EN)
    tb(s, Inches(3.0), Inches(2.2), Inches(9), Inches(1.0),
       title, 40, WHITE, True)
    line(s, Inches(3.0), Inches(3.5), Inches(3), color)
    tb(s, Inches(3.0), Inches(3.8), Inches(9), Inches(0.8),
       subtitle, 16, GRAY_LIGHT)
    footer(s)
    return s

def make_content(prs, title, color=ACCENT_BLUE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    line(s, Inches(0), Inches(0), SLIDE_W, color)
    tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
       title, 24, WHITE, True)
    footer(s)
    return s


# ═══════════════════════════════════════════════════════════════
# BUILD SLIDES
# ═══════════════════════════════════════════════════════════════

def create_ppt():
    prs = new_prs()

    # ── 1. 封面 ─────────────────────────────────────────────
    make_title_slide(prs,
        'BTC 多空強度預測指標',
        '系統架構與數據流程全解析  |  Dual-Model v7',
        '29 API Endpoints  |  130+ Features  |  XGBoost Dual-Model  |  4h Prediction')

    # ── 2. 系統總覽 ─────────────────────────────────────────
    s = make_content(prs, '系統總覽 — 端到端數據管線', ACCENT_BLUE)

    # Flow boxes
    flow = [
        ("數據採集", "3 家數據商\n29 個 API 端點\n每小時自動拉取", ACCENT_BLUE),
        ("特徵工程", "12 個特徵群組\n130+ 工程特徵\n純滾動計算", ACCENT_GREEN),
        ("模型推論", "Dual XGBoost v7\n方向分類 + 幅度回歸\n獨立特徵集", ACCENT_PURPLE),
        ("信號生成", "方向 / 信心 / 強度\nBBP 確認閘門\n動態死區 + 遲滯", ACCENT_ORANGE),
        ("圖表輸出", "4 面板 PNG 圖表\nTelegram 推送\nREST API 服務", ACCENT_RED),
    ]
    for i, (name, desc, color) in enumerate(flow):
        x = Inches(0.3 + i * 2.6)
        card(s, x, Inches(1.2), Inches(2.3), Inches(2.8))
        line(s, x, Inches(1.2), Inches(2.3), color)
        tb(s, x + Inches(0.1), Inches(1.4), Inches(2.1), Inches(0.4),
           name, 16, color, True, PP_ALIGN.CENTER)
        tb(s, x + Inches(0.1), Inches(2.0), Inches(2.1), Inches(1.8),
           desc, 11, GRAY_LIGHT, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_shape(s, x + Inches(2.35), Inches(2.3), Inches(0.2), Inches(0.35), color)

    # KPI cards
    kpis = [
        ("29", "API 端點", ACCENT_BLUE),
        ("130+", "特徵數", ACCENT_GREEN),
        ("2", "模型數", ACCENT_PURPLE),
        ("4h", "預測週期", ACCENT_ORANGE),
        ("4", "圖表面板", ACCENT_RED),
        ("1h", "更新頻率", ACCENT_CYAN),
    ]
    for i, (val, label, color) in enumerate(kpis):
        x = Inches(0.3 + i * 2.15)
        card(s, x, Inches(4.4), Inches(1.95), Inches(2.3), BG_CARD2)
        tb(s, x, Inches(4.7), Inches(1.95), Inches(0.8),
           val, 36, color, True, PP_ALIGN.CENTER, FONT_EN)
        tb(s, x, Inches(5.6), Inches(1.95), Inches(0.5),
           label, 13, GRAY_LIGHT, align=PP_ALIGN.CENTER)

    # ── 3. 章節: 數據採集 ────────────────────────────────────
    make_section(prs, '01', '數據採集層',
                 '多源即時數據收集 — Binance / Coinglass / Deribit', ACCENT_BLUE)

    # ── 4. 數據源概覽 ────────────────────────────────────────
    s = make_content(prs, '三大數據源總覽', ACCENT_BLUE)

    # Binance
    card(s, Inches(0.4), Inches(1.1), Inches(4.0), Inches(5.7))
    line(s, Inches(0.4), Inches(1.1), Inches(4.0), ACCENT_BLUE)
    tf = tb(s, Inches(0.6), Inches(1.3), Inches(3.6), Inches(0.5),
            'Binance Futures', 18, ACCENT_BLUE, True)
    ap(tf, '3 個 REST API 端點', 12, GRAY_DIM)
    ap(tf, '', 6)
    ap(tf, '/fapi/v1/klines', 12, ACCENT_CYAN, True, Pt(6))
    ap(tf, '  1h K 線 OHLCV + taker 成交量', 10, GRAY_LIGHT)
    ap(tf, '  500 bars / 次 (約 20 天歷史)', 10, GRAY_LIGHT)
    ap(tf, '  衍生: log_return, realized_vol, taker_delta', 10, GRAY_LIGHT)
    ap(tf, '  return_lag_1~10, quote_vol_zscore', 10, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, '/fapi/v1/depth', 12, ACCENT_CYAN, True, Pt(6))
    ap(tf, '  Order book 快照 (L20)', 10, GRAY_LIGHT)
    ap(tf, '  depth_imbalance, spread_bps', 10, GRAY_LIGHT)
    ap(tf, '  near_imbalance, bid_ask_ratio', 10, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, '/fapi/v1/aggTrades', 12, ACCENT_CYAN, True, Pt(6))
    ap(tf, '  逐筆交易, 2h 回看, 自動分頁', 10, GRAY_LIGHT)
    ap(tf, '  $100K 門檻: 大單 vs 小單分離', 10, GRAY_LIGHT)
    ap(tf, '  agg_large_delta, large_ratio', 10, GRAY_LIGHT)
    ap(tf, '  imbalance_div, delta_frac', 10, GRAY_LIGHT)

    # Coinglass
    card(s, Inches(4.7), Inches(1.1), Inches(4.6), Inches(5.7))
    line(s, Inches(4.7), Inches(1.1), Inches(4.6), ACCENT_GREEN)
    tf = tb(s, Inches(4.9), Inches(1.3), Inches(4.2), Inches(0.5),
            'Coinglass API v4', 18, ACCENT_GREEN, True)
    ap(tf, '24 個端點 (15 時序 + 9 快照)', 12, GRAY_DIM)
    ap(tf, '', 4)
    ap(tf, '時序端點 (原生 1h, 500 bars):', 13, WHITE, True, Pt(6))
    rows = [
        ("Open Interest", "OI 單所 / 聚合 / 幣本位"),
        ("清算數據", "清算單所 / 清算聚合"),
        ("多空比", "頂部帳戶 / 全局帳戶 / 持倉比"),
        ("資金費率", "Funding Rate OHLC"),
        ("Taker 成交量", "Taker 買賣量"),
        ("CVD", "期貨 CVD 聚合 / 現貨 CVD 聚合"),
        ("溢價/保證金", "Coinbase 溢價 / Bitfinex 保證金"),
    ]
    for cat, desc in rows:
        ap(tf, f'  {cat}', 10, ACCENT_CYAN, True, Pt(2))
        ap(tf, f'    {desc}', 9, GRAY_LIGHT, False, Pt(0))
    ap(tf, '', 4)
    ap(tf, '快照端點 (即時數據):', 13, WHITE, True, Pt(6))
    snaps = ["Options 最大痛點 / OI / P-C 比", "ETF 資金流 / AUM",
             "期貨/現貨淨流入 (多週期)", "恐懼貪婪指數 / HL 鯨魚倉位"]
    for snap in snaps:
        ap(tf, f'  {snap}', 9, GRAY_LIGHT, False, Pt(1))

    # Deribit + Quality
    card(s, Inches(9.6), Inches(1.1), Inches(3.4), Inches(3.0))
    line(s, Inches(9.6), Inches(1.1), Inches(3.4), ACCENT_PURPLE)
    tf = tb(s, Inches(9.8), Inches(1.3), Inches(3.0), Inches(0.5),
            'Deribit (免費)', 18, ACCENT_PURPLE, True)
    ap(tf, '2 個公開端點', 12, GRAY_DIM)
    ap(tf, '', 4)
    ap(tf, 'DVOL 波動率指數', 11, ACCENT_CYAN, True, Pt(4))
    ap(tf, '  bvol OHLC, intra_range', 9, GRAY_LIGHT)
    ap(tf, '  change_1h, zscore stubs', 9, GRAY_LIGHT)
    ap(tf, '', 3)
    ap(tf, 'Options Summary', 11, ACCENT_CYAN, True, Pt(4))
    ap(tf, '  P/C volume, IV skew', 9, GRAY_LIGHT)
    ap(tf, '  OTM put/call mean IV', 9, GRAY_LIGHT)

    card(s, Inches(9.6), Inches(4.4), Inches(3.4), Inches(2.4))
    line(s, Inches(9.6), Inches(4.4), Inches(3.4), ACCENT_YELLOW)
    tf = tb(s, Inches(9.8), Inches(4.55), Inches(3.0), Inches(0.4),
            '數據品質保障', 15, ACCENT_YELLOW, True)
    ap(tf, '', 4)
    ap(tf, '重試機制', 11, WHITE, True, Pt(4))
    ap(tf, '  指數退避 3 次 (2/4/8 秒)', 9, GRAY_LIGHT)
    ap(tf, '快取回退', 11, WHITE, True, Pt(4))
    ap(tf, '  .data_cache/ parquet', 9, GRAY_LIGHT)
    ap(tf, '新鮮度監控', 11, WHITE, True, Pt(4))
    ap(tf, '  > 3h 過時觸發 TG 告警', 9, GRAY_LIGHT)

    # ── 5. 章節: 特徵工程 ────────────────────────────────────
    make_section(prs, '02', '特徵工程',
                 '12 個特徵群組 | 130+ 工程特徵 | 純滾動計算 (無前視偏差)', ACCENT_GREEN)

    # ── 6. 特徵群組總覽 ──────────────────────────────────────
    s = make_content(prs, '特徵群組總覽 (130+ 個特徵)', ACCENT_GREEN)

    groups = [
        ("K 線衍生", "~33", ACCENT_BLUE,
         "log_return, realized_vol\ntaker_delta_ratio/ma/std\nreturn_lag_1~10\nquote_vol_zscore/ratio"),
        ("Coinglass", "~50+", ACCENT_GREEN,
         "OI delta/accel/pctchg\nliq surge/cascade\nL/S ratio, funding\ntaker delta/ratio"),
        ("Z-score", "~14", ACCENT_CYAN,
         "所有主要指標的\n24h 滾動標準化\n(zscore_win=24)"),
        ("交叉特徵", "~8", ACCENT_CYAN,
         "liq_x_oi, crowding\nconviction\ncvd_divergence\nmargin_funding_align"),
        ("方向特徵", "~40", ACCENT_PURPLE,
         "6 個子群組\n(見下頁詳細)"),
        ("波動率 BVOL", "~13", ACCENT_PURPLE,
         "Deribit DVOL OHLC\nintra_range, change_1h\n+ 滾動特徵佔位"),
        ("大單流 aggTrades", "~11", ACCENT_ORANGE,
         "agg_large_delta\nlarge_ratio, buy_ratio\nimbalance_div\ndelta_frac"),
        ("動能/時間/深度", "~12", ACCENT_ORANGE,
         "vol_acceleration\nentropy, squeeze_proxy\nhour/weekday cyclical\ndepth_imbalance"),
    ]

    for i, (name, count, color, desc) in enumerate(groups):
        col = i % 4
        row = i // 4
        x = Inches(0.3 + col * 3.2)
        y = Inches(1.1 + row * 3.05)
        card(s, x, y, Inches(3.0), Inches(2.8))
        line(s, x, y, Inches(3.0), color)
        tb(s, x + Inches(0.15), y + Inches(0.2), Inches(2.7), Inches(0.4),
           name, 14, color, True)
        tb(s, x + Inches(0.15), y + Inches(0.6), Inches(2.7), Inches(0.5),
           count + ' 個特徵', 24, WHITE, True)
        tb(s, x + Inches(0.15), y + Inches(1.3), Inches(2.7), Inches(1.3),
           desc, 10, GRAY_LIGHT)

    # ── 7. 方向特徵詳解 ──────────────────────────────────────
    s = make_content(prs, '方向特徵 — 6 個子群組 (40+ 個特徵)', ACCENT_PURPLE)

    dir_groups = [
        ("1. 大單流分離", "8~12", ACCENT_BLUE,
         "large_delta / small_delta\nlarge_buy_ratio / buy_bias\nlarge_delta_zscore\nma_4, slope_4, persistence\n\n原理: 大單方向 = 機構意圖"),
        ("2. 失衡持續性", "8", ACCENT_CYAN,
         "imb_1b / 3b / 5b / 8b\nimb_std_5b (低=持續)\nimb_sign_persistence\nimb_slope_4b, 3b_zscore\n\n原理: 單根失衡=雜訊, 持續=信號"),
        ("3. 吸收代理", "8", ACCENT_GREEN,
         "absorption_buy / sell / net\nabsorption_net_zscore\nbuy/sell_flow_price_impact\nimpact_asymmetry\n\n原理: 大流量+小波動=限價單吸收"),
        ("4. 短期動量", "8", ACCENT_ORANGE,
         "ret_1b / 2b / 3b / 5b\nwick_asymmetry\nbody_ratio / signed_body\nreversal_1v3, persistence\n\n原理: 價格行為方向性信號"),
        ("5. 情緒/反向", "5~7", ACCENT_RED,
         "funding_zscore / extreme\nls_ratio_zscore\ncrowding_zscore\nliq_imbalance_zscore/slope\noi_price_confirm\n\n原理: 極端倉位 = 均值回歸"),
        ("6. Order Book", "4~6", ACCENT_PURPLE,
         "obi_l1 / obi_l5 (+ zscore)\nobi_change\nobi_persistence\n\n原理: 掛單不平衡\n= 短期方向壓力"),
    ]

    for i, (name, count, color, desc) in enumerate(dir_groups):
        col = i % 3
        row = i // 3
        x = Inches(0.3 + col * 4.25)
        y = Inches(1.1 + row * 3.15)
        card(s, x, y, Inches(4.05), Inches(2.95))
        line(s, x, y, Inches(4.05), color)
        tf = tb(s, x + Inches(0.15), y + Inches(0.15), Inches(3.7), Inches(0.3),
                name, 13, color, True)
        ap(tf, count + ' 個特徵', 10, GRAY_DIM)
        tb(s, x + Inches(0.15), y + Inches(0.7), Inches(3.7), Inches(2.0),
           desc, 9, GRAY_LIGHT)

    # ── 8. 特徵計算流程 ──────────────────────────────────────
    s = make_content(prs, '特徵計算管線 — build_live_features()', ACCENT_GREEN)

    steps = [
        ("1", "原始數據\n輸入", "klines (500 bars)\ncg_data (15 端點)\ndepth (快照)\naggtrades (快照)\noptions_data (DVOL)", ACCENT_BLUE),
        ("2", "K 線衍生\n計算", "log_return, vol\ntaker_delta\nreturn lags\nquote_vol stats", ACCENT_GREEN),
        ("3", "Coinglass\n注入", "merge_asof()\n原生 1h 精確對齊\n15 端點全部注入\nz-score + 交叉特徵", ACCENT_GREEN),
        ("4", "動能/動態\n計算", "slope, momentum\nOI-price divergence\nvol dynamics\nabsorption / liq enhanced", ACCENT_PURPLE),
        ("5", "快照特徵\n注入", "BVOL (最後一根)\naggTrades (最後一根)\ndepth (最後一根)\n方向特徵 (全量)", ACCENT_ORANGE),
    ]
    for i, (num, name, desc, color) in enumerate(steps):
        x = Inches(0.3 + i * 2.6)
        card(s, x, Inches(1.1), Inches(2.35), Inches(4.2))
        line(s, x, Inches(1.1), Inches(2.35), color)
        tb(s, x, Inches(1.3), Inches(2.35), Inches(0.5),
           num, 28, color, True, PP_ALIGN.CENTER, FONT_EN)
        tb(s, x + Inches(0.1), Inches(1.8), Inches(2.15), Inches(0.6),
           name, 13, WHITE, True, PP_ALIGN.CENTER)
        tb(s, x + Inches(0.1), Inches(2.5), Inches(2.15), Inches(2.5),
           desc, 10, GRAY_LIGHT, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_shape(s, x + Inches(2.4), Inches(2.8), Inches(0.18), Inches(0.3), color)

    card(s, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.1), BG_CARD2)
    tf = tb(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
            '關鍵原則', 15, ACCENT_YELLOW, True)
    ap(tf, '  所有特徵使用純滾動計算 (trailing-only), 嚴格無前視偏差  |  Coinglass 原生 1h 數據使用 merge_asof 精確時間對齊  |  快照數據只設定最後一根 bar', 10, GRAY_LIGHT)

    # ── 9. 章節: 模型推論 ────────────────────────────────────
    make_section(prs, '03', '模型架構',
                 'Dual-Model v7 — 方向分類器 + 幅度回歸器 (獨立管線)', ACCENT_PURPLE)

    # ── 10. 雙模型架構 ───────────────────────────────────────
    s = make_content(prs, '雙模型 v7 架構', ACCENT_PURPLE)

    # Direction model
    card(s, Inches(0.3), Inches(1.1), Inches(6.2), Inches(4.2))
    line(s, Inches(0.3), Inches(1.1), Inches(6.2), ACCENT_BLUE)
    tf = tb(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.5),
            '方向模型 Direction Model', 20, ACCENT_BLUE, True)
    ap(tf, '', 4)
    ap(tf, '模型類型    XGBClassifier (二分類)', 13, GRAY_LIGHT)
    ap(tf, '特徵數量    89 個精選特徵', 13, GRAY_LIGHT)
    ap(tf, '輸出          P(UP) in [0, 1]', 13, GRAY_LIGHT)
    ap(tf, '', 6)
    ap(tf, '決策邏輯:', 14, WHITE, True)
    ap(tf, '  P(UP) > 0.60                UP', 12, ACCENT_GREEN, True)
    ap(tf, '  P(UP) < 0.40                DOWN', 12, ACCENT_RED, True)
    ap(tf, '  0.40 <= P(UP) <= 0.60   NEUTRAL', 12, GRAY_DIM, True)
    ap(tf, '', 4)
    ap(tf, 'BBP 確認閘門:', 14, WHITE, True)
    ap(tf, '  方向必須與多空力量指標同向', 12, GRAY_LIGHT)
    ap(tf, '  衝突時 (|BBP| > 0.15) 降級為 NEUTRAL', 12, GRAY_LIGHT)

    # Magnitude model
    card(s, Inches(6.8), Inches(1.1), Inches(6.2), Inches(4.2))
    line(s, Inches(6.8), Inches(1.1), Inches(6.2), ACCENT_ORANGE)
    tf = tb(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.5),
            '幅度模型 Magnitude Model', 20, ACCENT_ORANGE, True)
    ap(tf, '', 4)
    ap(tf, '模型類型    XGBRegressor', 13, GRAY_LIGHT)
    ap(tf, '特徵數量    78 個精選特徵', 13, GRAY_LIGHT)
    ap(tf, '輸出          |return_4h| >= 0', 13, GRAY_LIGHT)
    ap(tf, '', 6)
    ap(tf, '組合輸出:', 14, WHITE, True)
    ap(tf, '  pred_return_4h = sign(direction) x magnitude', 12, ACCENT_CYAN, True)
    ap(tf, '', 4)
    ap(tf, '信心評分:', 14, WHITE, True)
    ap(tf, '  |mag_pred| 的擴展百分位排名 (0~100)', 12, GRAY_LIGHT)
    ap(tf, '  方向確信度加成: x (0.7 + 0.3 x |P(UP)-0.5| x 2)', 12, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, '強度等級:', 14, WHITE, True)
    ap(tf, '  >= 80 Strong  |  >= 65 Moderate  |  < 65 Weak', 12, GRAY_LIGHT)

    # Fallback
    card(s, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.1), BG_CARD2)
    tf = tb(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
            '模型回退鏈', 15, ACCENT_CYAN, True)
    ap(tf, '  Dual Model (v7)  ───>  Regime Models (v5/v6, per-target per-regime)  ───>  Legacy Single Model (v2)', 13, WHITE, True)
    ap(tf, '  自動偵測 model_artifacts/ 目錄下的模型檔案, 優先使用最新版本', 10, GRAY_DIM)

    # ── 11. 信號生成流程 ─────────────────────────────────────
    s = make_content(prs, '信號生成管線 (6 步驟)', ACCENT_PURPLE)

    steps = [
        ("1", "方向分類", "P(UP) 閾值判定\n> 0.60 = UP\n< 0.40 = DOWN\nelse NEUTRAL", ACCENT_BLUE),
        ("2", "幅度預測", "|return_4h|\n預測值 >= 0\n獨立於方向", ACCENT_GREEN),
        ("3", "信心評分", "|mag| 百分位\n+ 方向確信加成\n0~100 分\n需 30 bar 暖機", ACCENT_PURPLE),
        ("4", "強度分級", ">= 80 Strong\n>= 65 Moderate\n< 65 Weak", ACCENT_ORANGE),
        ("5", "BBP 閘門", "方向須與\nBBP 同向\n衝突則降級\n為 NEUTRAL", ACCENT_CYAN),
        ("6", "Regime", "CHOPPY x1.6\nTREND x0.9\n遲滯 x1.4\n冷卻 >= 1 bar", ACCENT_RED),
    ]
    for i, (num, name, desc, color) in enumerate(steps):
        x = Inches(0.3 + i * 2.15)
        card(s, x, Inches(1.1), Inches(1.95), Inches(3.8))
        line(s, x, Inches(1.1), Inches(1.95), color)
        tb(s, x, Inches(1.3), Inches(1.95), Inches(0.5),
           num, 28, color, True, PP_ALIGN.CENTER, FONT_EN)
        tb(s, x + Inches(0.1), Inches(1.9), Inches(1.75), Inches(0.4),
           name, 14, WHITE, True, PP_ALIGN.CENTER)
        tb(s, x + Inches(0.1), Inches(2.5), Inches(1.75), Inches(2.0),
           desc, 10, GRAY_LIGHT, align=PP_ALIGN.CENTER)
        if i < 5:
            arrow_shape(s, x + Inches(2.0), Inches(2.7), Inches(0.13), Inches(0.25), color)

    # Regime box
    card(s, Inches(0.3), Inches(5.2), Inches(12.7), Inches(1.5))
    tf = tb(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
            'Regime 偵測 (市場狀態判定)', 14, ACCENT_RED, True)
    ap(tf, '', 3)
    ap(tf, '  WARMUP (前 168 bars / 7 天)       CHOPPY (vol <= 0.6 或 |24h ret| < 0.5%)       TRENDING_BULL / BEAR (vol > 0.6 且 |24h ret| > 0.5%)', 10, GRAY_LIGHT)
    ap(tf, '  遲滯機制: 翻轉方向需要 1.4 倍的反向信號強度  |  冷卻機制: 翻轉後至少持續 1 根 bar', 10, GRAY_DIM)

    # ── 12. 章節: 圖表輸出 ───────────────────────────────────
    make_section(prs, '04', '圖表與輸出',
                 '4 面板圖表 | Telegram 推送 | REST API | 數據持久化', ACCENT_ORANGE)

    # ── 13. 4 面板圖表 ───────────────────────────────────────
    s = make_content(prs, '4 面板圖表輸出', ACCENT_ORANGE)

    panels = [
        ("Panel 1", "信心熱力圖", "0.8x", ACCENT_PURPLE,
         "紫色漸層條\n0~100% 信心分數\n深色 = 高信心\n顯示最近 200 根 bar"),
        ("Panel 2", "K 線 + 方向三角形", "8x", ACCENT_GREEN,
         "綠色/紅色 K 線\n三角形 (Moderate/Strong)\nUP = 綠色向上三角\nDOWN = 紅色向下三角\nWeak 不顯示三角形"),
        ("Panel 3", "Magnitude 幅度", "2x", ACCENT_BLUE,
         "預測 |return_4h| (%)\nUP: 綠色柱 (零線上)\nDOWN: 紅色柱 (零線下)\nNEUTRAL: 灰色柱 (零線上)\n幅度獨立於方向顯示"),
        ("Panel 4", "Bull/Bear Power", "2x", ACCENT_RED,
         "BBP 合成指標 [-1, 1]\n綠色 = 多方力量\n紅色 = 空方力量\n由 5 個 CG 指標合成:\nOI/funding/taker/LS"),
    ]
    for i, (panel, name, ratio, color, desc) in enumerate(panels):
        x = Inches(0.3 + i * 3.2)
        card(s, x, Inches(1.1), Inches(3.0), Inches(5.6))
        line(s, x, Inches(1.1), Inches(3.0), color)
        tf = tb(s, x + Inches(0.1), Inches(1.3), Inches(2.8), Inches(0.3),
                panel, 11, GRAY_DIM)
        ap(tf, name, 15, color, True, Pt(2))
        tb(s, x + Inches(0.1), Inches(2.1), Inches(2.8), Inches(0.3),
           f'高度比例: {ratio}', 10, GRAY_DIM)
        tb(s, x + Inches(0.1), Inches(2.5), Inches(2.8), Inches(3.8),
           desc, 10, GRAY_LIGHT)

    # ── 14. 輸出管道 ─────────────────────────────────────────
    s = make_content(prs, '輸出與傳遞管道', ACCENT_ORANGE)

    # Telegram
    card(s, Inches(0.3), Inches(1.1), Inches(4.2), Inches(5.6))
    line(s, Inches(0.3), Inches(1.1), Inches(4.2), ACCENT_CYAN)
    tf = tb(s, Inches(0.5), Inches(1.3), Inches(3.8), Inches(0.4),
            'Telegram Bot 推送', 18, ACCENT_CYAN, True)
    ap(tf, '', 4)
    ap(tf, '圖表傳送:', 13, WHITE, True, Pt(6))
    ap(tf, '  4 面板 PNG + 說明文字', 11, GRAY_LIGHT)
    ap(tf, '  方向 / 信心 / 強度 / regime', 11, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, '快速鍵盤 (8 個按鈕):', 13, WHITE, True, Pt(4))
    ap(tf, '  Chart | Perf | DB | Flow BTC', 11, GRAY_LIGHT)
    ap(tf, '  Flow All | Status | Events | Help', 11, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, '自動告警:', 13, WHITE, True, Pt(4))
    ap(tf, '  數據新鮮度超時 (> 3h)', 11, GRAY_LIGHT)
    ap(tf, '  全端點失效告警', 11, GRAY_LIGHT)
    ap(tf, '  Strong 信號特殊提醒', 11, GRAY_LIGHT)

    # REST API
    card(s, Inches(4.8), Inches(1.1), Inches(4.2), Inches(5.6))
    line(s, Inches(4.8), Inches(1.1), Inches(4.2), ACCENT_GREEN)
    tf = tb(s, Inches(5.0), Inches(1.3), Inches(3.8), Inches(0.4),
            'REST API (Flask)', 18, ACCENT_GREEN, True)
    ap(tf, '', 4)
    routes = [
        ("GET /", "最新圖表 PNG"),
        ("GET /health", "系統健康檢查"),
        ("GET /json", "最新預測 JSON"),
        ("POST /force-update", "手動觸發更新"),
        ("GET /indicator-status", "詳細狀態 + CG 健康"),
        ("GET /db-diag", "MySQL 診斷"),
        ("GET /indicator-perf", "滾動準確率指標"),
        ("GET /diag", "完整診斷頁面"),
    ]
    for route, desc in routes:
        ap(tf, '', 2)
        ap(tf, f'  {route}', 10, ACCENT_CYAN, True, Pt(1))
        ap(tf, f'    {desc}', 9, GRAY_LIGHT, False, Pt(0))

    # Persistence
    card(s, Inches(9.3), Inches(1.1), Inches(3.7), Inches(5.6))
    line(s, Inches(9.3), Inches(1.1), Inches(3.7), ACCENT_PURPLE)
    tf = tb(s, Inches(9.5), Inches(1.3), Inches(3.3), Inches(0.4),
            '數據持久化', 18, ACCENT_PURPLE, True)
    ap(tf, '', 4)
    ap(tf, 'MySQL', 14, WHITE, True, Pt(6))
    ap(tf, '  indicator_history 表', 11, GRAY_LIGHT)
    ap(tf, '  14+ 欄位, UPSERT', 11, GRAY_LIGHT)
    ap(tf, '  每根 bar 自動寫入', 11, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, 'Parquet', 14, WHITE, True, Pt(6))
    ap(tf, '  完整歷史紀錄', 11, GRAY_LIGHT)
    ap(tf, '  重啟自動恢復', 11, GRAY_LIGHT)
    ap(tf, '  Git 版本控制', 11, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, '.data_cache/', 14, WHITE, True, Pt(6))
    ap(tf, '  API 回退快取', 11, GRAY_LIGHT)
    ap(tf, '  per-source parquet', 11, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, 'Snapshots', 14, WHITE, True, Pt(6))
    ap(tf, '  depth / aggTrades', 11, GRAY_LIGHT)
    ap(tf, '  options / sentiment', 11, GRAY_LIGHT)
    ap(tf, '  存入 MySQL', 11, GRAY_LIGHT)

    # ── 15. 章節: 基礎設施 ───────────────────────────────────
    make_section(prs, '05', '基礎設施',
                 'Railway 部署 | MySQL 8.0 | Flask + APScheduler | 執行緒模型', ACCENT_RED)

    # ── 16. 技術棧 ───────────────────────────────────────────
    s = make_content(prs, '技術棧與部署架構', ACCENT_RED)

    # Tech stack
    card(s, Inches(0.3), Inches(1.1), Inches(5.5), Inches(5.6))
    line(s, Inches(0.3), Inches(1.1), Inches(5.5), ACCENT_BLUE)
    tf = tb(s, Inches(0.5), Inches(1.3), Inches(5.1), Inches(0.4),
            '技術棧', 18, ACCENT_BLUE, True)
    ap(tf, '', 4)
    stack = [
        ("執行環境", "Python 3.11"),
        ("機器學習", "XGBoost (Classifier + Regressor)"),
        ("數據處理", "Pandas, NumPy, SciPy"),
        ("Web 框架", "Flask"),
        ("資料庫", "MySQL 8.0 (Railway 託管)"),
        ("部署平台", "Railway (git push 自動部署)"),
        ("排程器", "APScheduler (1h 間隔)"),
        ("圖表引擎", "Matplotlib (Agg backend)"),
        ("推送服務", "Telegram Bot API"),
        ("環境配置", "python-dotenv (.env)"),
        ("HTTP 客戶端", "Requests (重試 + 退避)"),
    ]
    for comp, tech in stack:
        ap(tf, f'  {comp}', 11, ACCENT_CYAN, True, Pt(3))
        ap(tf, f'    {tech}', 10, GRAY_LIGHT, False, Pt(0))

    # Threading + Deploy
    card(s, Inches(6.1), Inches(1.1), Inches(6.9), Inches(3.0))
    line(s, Inches(6.1), Inches(1.1), Inches(6.9), ACCENT_ORANGE)
    tf = tb(s, Inches(6.3), Inches(1.3), Inches(6.5), Inches(0.4),
            '執行緒模型', 18, ACCENT_ORANGE, True)
    ap(tf, '', 4)
    ap(tf, 'Flask 主執行緒 (port $PORT)', 13, WHITE, True, Pt(4))
    ap(tf, '  HTTP 請求處理 + threading.Lock 保護共享狀態', 10, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, 'APScheduler 背景執行緒 (每小時觸發)', 13, WHITE, True, Pt(4))
    ap(tf, '  1. 拉取數據 (Binance + CG + Deribit)    5. 渲染 4 面板圖表', 10, GRAY_LIGHT)
    ap(tf, '  2. 建構特徵 (130+ 欄位)                     6. 儲存 MySQL + Parquet', 10, GRAY_LIGHT)
    ap(tf, '  3. 執行雙模型推論                               7. 推送 Telegram 圖片', 10, GRAY_LIGHT)
    ap(tf, '  4. 回填歷史 mag_pred (如需要)', 10, GRAY_LIGHT)

    card(s, Inches(6.1), Inches(4.4), Inches(6.9), Inches(2.3))
    line(s, Inches(6.1), Inches(4.4), Inches(6.9), ACCENT_RED)
    tf = tb(s, Inches(6.3), Inches(4.6), Inches(6.5), Inches(0.4),
            'Railway 部署', 18, ACCENT_RED, True)
    ap(tf, '', 4)
    ap(tf, '  git push main 觸發自動部署', 12, GRAY_LIGHT)
    ap(tf, '  Service 1: 主 Bot (Dockerfile)', 12, GRAY_LIGHT)
    ap(tf, '  Service 2: 市場數據 (Dockerfile.marketdata)', 12, GRAY_LIGHT)
    ap(tf, '  MySQL: Railway 內部 mysql.railway.internal', 12, GRAY_LIGHT)
    ap(tf, '  環境變數: Railway 設定 > .env 檔案', 12, GRAY_LIGHT)

    # ── 17. 檔案結構 ─────────────────────────────────────────
    s = make_content(prs, '核心檔案結構', ACCENT_RED)

    card(s, Inches(0.3), Inches(1.1), Inches(7.5), Inches(5.6))
    line(s, Inches(0.3), Inches(1.1), Inches(7.5), ACCENT_BLUE)
    tf = tb(s, Inches(0.5), Inches(1.3), Inches(7.1), Inches(0.4),
            'indicator/ (生產模組)', 16, ACCENT_BLUE, True)
    files = [
        ("app.py", "900+", "Flask API, Telegram, 更新排程, 主入口"),
        ("inference.py", "700+", "IndicatorEngine v7: 模型載入, 推論, regime 偵測"),
        ("data_fetcher.py", "730+", "13 個 fetch 函數, 重試/快取, 29 API 端點"),
        ("feature_builder_live.py", "530+", "特徵工程: 130+ 欄位從原始 API 計算"),
        ("chart_renderer.py", "300+", "Matplotlib 4 面板圖表渲染 (PNG)"),
        ("auto_update.py", "400+", "獨立排程器, 歷史持久化, TG 告警"),
        ("snapshot_collector.py", "350+", "Depth/aggTrades/options/sentiment 快照"),
    ]
    for fname, loc, purpose in files:
        ap(tf, '', 2)
        ap(tf, f'  {fname}', 11, ACCENT_CYAN, True, Pt(2))
        ap(tf, f'    {loc} 行  |  {purpose}', 9, GRAY_LIGHT, False, Pt(0))

    card(s, Inches(0.3), Inches(1.1 + 4.0), Inches(7.5), Inches(1.6))
    tf = tb(s, Inches(0.5), Inches(5.3), Inches(7.1), Inches(0.4),
            'research/ (訓練與實驗)', 16, ACCENT_PURPLE, True)
    ap(tf, '', 3)
    ap(tf, '  direction_features.py    450+ 行  |  6 個子群組, 40+ 方向特徵', 10, GRAY_LIGHT, False, Pt(2))
    ap(tf, '  dual_model/train_*.py   300+ 行  |  XGB 分類器/回歸器訓練管線', 10, GRAY_LIGHT, False, Pt(1))
    ap(tf, '  dual_model/evaluate_*.py 200+ 行  |  OOS 評估, IC/ICIR, calibration', 10, GRAY_LIGHT, False, Pt(1))

    # Model artifacts
    card(s, Inches(8.1), Inches(1.1), Inches(4.9), Inches(5.6))
    line(s, Inches(8.1), Inches(1.1), Inches(4.9), ACCENT_PURPLE)
    tf = tb(s, Inches(8.3), Inches(1.3), Inches(4.5), Inches(0.4),
            'model_artifacts/', 16, ACCENT_PURPLE, True)
    ap(tf, '', 4)
    ap(tf, 'dual_model/ (v7 主要)', 13, WHITE, True, Pt(6))
    ap(tf, '  direction_xgb.json', 10, GRAY_LIGHT)
    ap(tf, '  direction_feature_cols.json (89)', 10, GRAY_LIGHT)
    ap(tf, '  magnitude_xgb.json', 10, GRAY_LIGHT)
    ap(tf, '  magnitude_feature_cols.json (78)', 10, GRAY_LIGHT)
    ap(tf, '  training_stats.json', 10, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, 'direction_models/ (v5 回退)', 13, WHITE, True, Pt(6))
    ap(tf, '  direction_xgb.json (170+ 特徵)', 10, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, 'regime_models/ (v5 回退)', 13, WHITE, True, Pt(6))
    ap(tf, '  per-target per-regime 模型', 10, GRAY_LIGHT)
    ap(tf, '', 4)
    ap(tf, '其他', 13, WHITE, True, Pt(6))
    ap(tf, '  indicator_history.parquet', 10, GRAY_LIGHT)
    ap(tf, '  xgb_model.json (v2 legacy)', 10, GRAY_LIGHT)
    ap(tf, '  .data_cache/ (API 快取)', 10, GRAY_LIGHT)

    # ── 18. 評估指標 ─────────────────────────────────────────
    s = make_content(prs, '模型評估指標', ACCENT_GREEN)

    card(s, Inches(0.3), Inches(1.1), Inches(12.7), Inches(1.5))
    tf = tb(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
            '核心定位', 18, ACCENT_YELLOW, True)
    ap(tf, '', 4)
    ap(tf, '  這是「多空強度預測指標」, 不是交易策略  |  所有評估以預測品質為準, 不做交易績效回測', 13, WHITE, True)
    ap(tf, '  嚴禁延伸到: entry/exit 規則, TP/SL 設計, 倉位管理, 自動下單, 策略回測', 11, ACCENT_RED)

    metrics = [
        ("Spearman IC", "> 0.05", "預測值與實際收益的排序相關係數", ACCENT_BLUE),
        ("ICIR", "> 0.3", "IC / std(IC), 衡量預測穩定性", ACCENT_BLUE),
        ("方向準確率", "> 58%", "pred_direction vs actual 實際方向匹配", ACCENT_GREEN),
        ("Calibration", "Monotonic", "預測越強 -> 實際收益越高 (單調遞增)", ACCENT_PURPLE),
        ("信心分布", "均勻", "不應全部集中在 Weak, 需合理分散", ACCENT_ORANGE),
    ]
    for i, (name, target, desc, color) in enumerate(metrics):
        y = Inches(2.9 + i * 0.85)
        card(s, Inches(0.3), y, Inches(12.7), Inches(0.72))
        tb(s, Inches(0.6), y + Inches(0.08), Inches(2.5), Inches(0.3),
           name, 14, color, True)
        tb(s, Inches(3.3), y + Inches(0.08), Inches(1.5), Inches(0.3),
           target, 14, WHITE, True, font=FONT_EN)
        tb(s, Inches(5.0), y + Inches(0.08), Inches(7.5), Inches(0.3),
           desc, 12, GRAY_LIGHT)

    # ── 19. 總結 ─────────────────────────────────────────────
    s = make_content(prs, '系統架構總結', ACCENT_BLUE)

    summary = [
        ("數據來源", "3 家數據商 (Binance, Coinglass, Deribit), 29 個 API 端點", ACCENT_BLUE),
        ("特徵工程", "12 個群組, 130+ 工程特徵, 純滾動計算, 無前視偏差", ACCENT_GREEN),
        ("模型架構", "Dual XGBoost v7: 方向分類 (89 feat) + 幅度回歸 (78 feat)", ACCENT_PURPLE),
        ("信號管線", "6 步驟: 方向 -> 幅度 -> 信心 -> 強度 -> BBP 閘門 -> Regime", ACCENT_PURPLE),
        ("圖表輸出", "4 面板 (信心熱力圖 + K 線 + 幅度柱 + BBP), 200 bars, UTC+8", ACCENT_ORANGE),
        ("傳遞管道", "Telegram (圖片 + 8 鍵) + REST API (10 路由) + MySQL + Parquet", ACCENT_ORANGE),
        ("部署架構", "Railway 自動部署, Flask + APScheduler, 1h 更新週期", ACCENT_RED),
        ("品質保障", "重試/快取/新鮮度監控, BBP 確認閘門, 動態死區, 遲滯機制", ACCENT_RED),
    ]
    for i, (key, val, color) in enumerate(summary):
        y = Inches(1.1 + i * 0.75)
        card(s, Inches(0.3), y, Inches(12.7), Inches(0.62))
        tb(s, Inches(0.6), y + Inches(0.05), Inches(2.2), Inches(0.4),
           key, 14, color, True)
        tb(s, Inches(3.0), y + Inches(0.05), Inches(9.7), Inches(0.4),
           val, 12, GRAY_LIGHT)

    # Save
    out = 'BTC_系統架構_v7.pptx'
    prs.save(out)
    print(f'OK: {out} ({len(prs.slides)} slides)')


if __name__ == '__main__':
    create_ppt()
