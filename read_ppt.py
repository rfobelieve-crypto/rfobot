import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
prs = Presentation('BTC_Market_Intelligence_Indicator.pptx')
for i, slide in enumerate(prs.slides):
    print(f'=== Slide {i+1} ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                print(para.text)
    print()
